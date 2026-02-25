"""
PPT转视频转换器 - 基于pptx2mp4原理实现
使用python-pptx + PIL + FFmpeg，无需Office软件
"""

import os
import io
import time
import shutil
import subprocess
import tempfile
from typing import Dict, Any
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from .base import BaseConverter


class PptToVideoPptx2mp4Converter(BaseConverter):
    """PPT转视频转换器（基于pptx2mp4原理）"""
    
    def __init__(self):
        super().__init__()
        self.ffmpeg_path = self._find_ffmpeg()
        
    def _find_ffmpeg(self) -> str:
        """查找FFmpeg"""
        ffmpeg = shutil.which('ffmpeg')
        if ffmpeg:
            return ffmpeg
        common_paths = [
            r'C:\ffmpeg\bin\ffmpeg.exe',
            r'C:\Program Files\ffmpeg\bin\ffmpeg.exe',
        ]
        for path in common_paths:
            if os.path.exists(path):
                return path
        return None
    
    def convert(self, input_path: str, output_path: str, **options) -> Dict[str, Any]:
        """转换PPT为视频"""
        
        print(f"[Pptx2Mp4] 开始转换: {input_path}")
        
        self.validate_input(input_path)
        self.update_progress(input_path, 5)
        
        if not self.ffmpeg_path:
            raise Exception("FFmpeg未安装。请安装FFmpeg")
        
        temp_dir = tempfile.mkdtemp(prefix='ppt_video_')
        
        try:
            # 步骤1：解析PPT并生成图片
            print(f"[Pptx2Mp4] 步骤1: 解析PPT...")
            self.update_progress(input_path, 20)
            
            image_files = self._ppt_to_images(input_path, temp_dir)
            
            if not image_files:
                raise Exception("未能生成图片")
            
            print(f"[Pptx2Mp4] 生成了 {len(image_files)} 张图片")
            self.update_progress(input_path, 60)
            
            # 步骤2：合成视频
            print(f"[Pptx2Mp4] 步骤2: 合成视频...")
            
            slide_duration = options.get('slide_duration', 3)
            fps = options.get('fps', 24)
            resolution = options.get('resolution', 720)
            
            self._create_video(image_files, output_path, slide_duration, fps, resolution)
            
            self.update_progress(input_path, 100)
            print(f"[Pptx2Mp4] 转换完成: {output_path}")
            
            return {
                'success': True,
                'output_path': output_path,
                'size': os.path.getsize(output_path),
                'slides': len(image_files),
                'duration': len(image_files) * slide_duration,
                'converter': 'pptx2mp4'
            }
            
        except Exception as e:
            raise Exception(f"PPT转视频失败: {str(e)}")
            
        finally:
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
    
    def _ppt_to_images(self, ppt_path: str, output_dir: str) -> list:
        """将PPT转换为图片"""
        
        try:
            print(f"[Pptx2Mp4] 正在打开PPT文件...")
            # 打开PPT
            prs = Presentation(ppt_path)
            slide_count = len(prs.slides)
            print(f"[Pptx2Mp4] PPT包含 {slide_count} 页")
            
            # 获取PPT尺寸
            ppt_width = prs.slide_width
            ppt_height = prs.slide_height
            
            # 转换为像素 (EMU to pixels, 1 inch = 914400 EMU, 1 inch = 96 pixels)
            width_px = int(ppt_width / 914400 * 96)
            height_px = int(ppt_height / 914400 * 96)
            
            # 限制最大尺寸
            max_width = 1920
            max_height = 1080
            if width_px > max_width or height_px > max_height:
                ratio = min(max_width / width_px, max_height / height_px)
                width_px = int(width_px * ratio)
                height_px = int(height_px * ratio)
            
            print(f"[Pptx2Mp4] 图片尺寸: {width_px}x{height_px}")
            
            image_files = []
            
            for i, slide in enumerate(prs.slides, 1):
                print(f"[Pptx2Mp4] 正在处理第 {i}/{slide_count} 页...")
                output_file = os.path.join(output_dir, f"slide_{i:04d}.png")
                
                # 创建白色背景图片
                img = Image.new('RGB', (width_px, height_px), 'white')
                draw = ImageDraw.Draw(img)
                
                # 渲染幻灯片内容（简化版）
                try:
                    self._render_slide_simple(slide, img, draw, width_px, height_px, ppt_width, ppt_height)
                except Exception as e:
                    print(f"[Pptx2Mp4] 渲染第{i}页时出错: {e}")
                
                # 保存图片
                img.save(output_file, 'PNG')
                image_files.append(output_file)
                
                if i % 5 == 0 or i == slide_count:
                    print(f"[Pptx2Mp4] 已生成 {i}/{slide_count}")
            
            return image_files
            
        except Exception as e:
            print(f"[Pptx2Mp4] 解析PPT失败: {e}")
            raise
    
    def _render_slide_simple(self, slide, img, draw, width_px, height_px, ppt_width, ppt_height):
        """简化版幻灯片渲染 - 只渲染文本"""
        
        # 只处理文本框
        y_offset = 50
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text[:200]  # 限制文本长度
                    
                    # 使用系统字体
                    try:
                        font = ImageFont.truetype("msyh.ttc", 24)
                    except:
                        try:
                            font = ImageFont.truetype("arial.ttf", 20)
                        except:
                            font = ImageFont.load_default()
                    
                    # 绘制文本
                    draw.text((50, y_offset), text, fill=(0, 0, 0), font=font)
                    y_offset += 40
            except:
                continue
    
    def _create_video(self, image_files: list, output_path: str,
                     slide_duration: float, fps: int, resolution: int):
        """使用FFmpeg创建视频"""
        
        if not image_files:
            raise Exception("没有图片")
        
        list_file = os.path.join(os.path.dirname(image_files[0]), 'filelist.txt')
        
        with open(list_file, 'w', encoding='utf-8') as f:
            for img in image_files:
                img_path = os.path.abspath(img).replace('\\', '/')
                f.write(f"file '{img_path}'\n")
                f.write(f"duration {slide_duration}\n")
            img_path = os.path.abspath(image_files[-1]).replace('\\', '/')
            f.write(f"file '{img_path}'\n")
        
        cmd = [
            self.ffmpeg_path,
            '-f', 'concat',
            '-safe', '0',
            '-i', list_file,
            '-vf', f'scale=-2:{resolution}',
            '-c:v', 'libx264',
            '-preset', 'fast',
            '-crf', '23',
            '-pix_fmt', 'yuv420p',
            '-r', str(fps),
            '-y',
            output_path
        ]
        
        print(f"[Pptx2Mp4] 执行FFmpeg...")
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        
        if result.returncode != 0:
            raise Exception(f"FFmpeg失败: {result.stderr}")
        
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            raise Exception("视频文件未生成")
        
        print(f"[Pptx2Mp4] 视频已创建")
