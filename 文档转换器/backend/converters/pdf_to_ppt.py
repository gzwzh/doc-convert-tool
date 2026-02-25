import os
import tempfile
import io
from pathlib import Path
from PIL import Image
from .base import BaseConverter
from typing import Dict, Any


class PdfToPptConverter(BaseConverter):
    """PDF 到 PPT 转换器（优化版 - 参考 conversion_core）
    
    优化内容：
    1. 添加进度回调支持
    2. 多策略降级（PyMuPDF图片模式 → pdfplumber文本模式）
    3. 高质量页面渲染
    4. 自动适应幻灯片尺寸
    5. 保持宽高比
    
    实现方式：
    - 优先使用 PyMuPDF 将每页 PDF 渲染为高质量图片，插入到 PPT 幻灯片
    - 降级方案使用 pdfplumber 提取文本内容创建文本幻灯片
    """
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['pptx']
    
    def convert(self, input_path: str, output_path: str, **options) -> Dict[str, Any]:
        """将 PDF 转换为 PPT（多策略）"""
        try:
            self.validate_input(input_path)
            self.update_progress(input_path, 5)
            
            # 检查依赖库
            try:
                import fitz  # PyMuPDF
                pymupdf_available = True
                print(f"[PDF转PPT] PyMuPDF 可用，版本: {fitz.version}")
            except ImportError as e:
                pymupdf_available = False
                print(f"[PDF转PPT] PyMuPDF 不可用: {e}，将使用 pdfplumber（仅文本模式）")
            
            try:
                import pdfplumber
            except ImportError as e:
                if not pymupdf_available:
                    raise Exception(f"pdfplumber 库导入失败: {str(e)}")
            
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt, Emu
            except ImportError as e:
                raise Exception(f"python-pptx 库导入失败: {str(e)}")
            
            self.update_progress(input_path, 15)
            
            # 获取选项
            quality = options.get('quality', 'high')  # 'ultra', 'high', 'medium'
            slide_width = options.get('slide_width', 10)  # 英寸
            slide_height = options.get('slide_height', 7.5)  # 英寸
            
            # 解析页面范围
            import fitz
            with fitz.open(input_path) as doc:
                total_pages = doc.page_count
                
            raw_page_range = options.get('pdf_page_range') or options.get('page_range')
            pages = self.parse_page_range(raw_page_range, total_pages=total_pages)
            
            # 确定要处理的页面
            if pages is None:
                pages_to_process = range(total_pages)
            else:
                pages_to_process = pages
            
            # 质量设置
            quality_map = {
                'ultra': 3.0,  # 3x 缩放
                'high': 2.0,   # 2x 缩放
                'medium': 1.5  # 1.5x 缩放
            }
            scale_factor = quality_map.get(quality, 2.0)
            
            # 创建 PPT 演示文稿
            prs = Presentation()
            prs.slide_width = Inches(slide_width)
            prs.slide_height = Inches(slide_height)
            
            # 策略1: 使用 PyMuPDF（推荐，支持图片和完整页面渲染）
            if pymupdf_available:
                result = self._convert_with_pymupdf(
                    input_path, output_path, prs, 
                    pages=pages_to_process,
                    scale_factor=scale_factor
                )
                return result
            else:
                # 策略2: 使用 pdfplumber（降级方案，仅文本）
                print("[PDF转PPT] 使用 pdfplumber 降级方案（仅文本）")
                self.update_progress(input_path, 20)
                result = self._convert_with_pdfplumber(
                    input_path, output_path, prs,
                    pages=pages_to_process
                )
                result['method'] = 'pdfplumber_fallback'
                result['warning'] = 'PyMuPDF not available, used text-only conversion'
                return result
            
        except Exception as e:
            self.cleanup_on_error(output_path)
            raise Exception(f"PDF to PPT conversion failed: {str(e)}")
    
    def _convert_with_pymupdf(self, input_path: str, output_path: str, 
                              prs, pages=None, scale_factor: float = 2.0) -> Dict[str, Any]:
        """使用 PyMuPDF 转换（主策略 - 高质量图片模式）"""
        import fitz
        from pptx.util import Inches, Emu
        
        doc = fitz.open(input_path)
        total_pages = len(doc)
        
        if total_pages == 0:
            doc.close()
            raise Exception("PDF 文件为空，没有可转换的页面")
        
        # 确定要处理的页面
        if pages is None:
            pages_to_process = range(total_pages)
        else:
            pages_to_process = pages
            
        self.update_progress(input_path, 25)
        
        processed_count = 0
        total_to_process = len(pages_to_process)
        
        for page_idx in pages_to_process:
            if page_idx < 0 or page_idx >= total_pages:
                continue
                
            page = doc[page_idx]
            
            # 创建新幻灯片（空白布局）
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 提取页面为高质量图片
            mat = fitz.Matrix(scale_factor, scale_factor)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            
            # 保存为临时图片
            img_data = pix.tobytes("png")
            
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                tmp_file.write(img_data)
            
            try:
                # 计算图片在幻灯片中的位置和大小（保持宽高比）
                img = Image.open(tmp_path)
                img_width, img_height = img.size
                img.close()
                
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # 计算缩放比例（留出边距）
                margin = Inches(0.3)
                available_width = slide_width - 2 * margin
                available_height = slide_height - 2 * margin
                
                # 将图片尺寸转换为 EMU（English Metric Units）
                # 假设 96 DPI
                img_width_emu = Emu(img_width * 914400 / 96)
                img_height_emu = Emu(img_height * 914400 / 96)
                
                # 计算缩放比例
                width_ratio = available_width / img_width_emu
                height_ratio = available_height / img_height_emu
                scale = min(width_ratio, height_ratio, 1.0)  # 不放大，只缩小
                
                pic_width = int(img_width_emu * scale)
                pic_height = int(img_height_emu * scale)
                
                # 居中放置
                left = int((slide_width - pic_width) / 2)
                top = int((slide_height - pic_height) / 2)
                
                # 添加图片到幻灯片
                slide.shapes.add_picture(tmp_path, left, top, pic_width, pic_height)
                
            finally:
                # 删除临时文件
                try:
                    os.unlink(tmp_path)
                except:
                    pass
            
            # 更新进度 (25% ~ 90%)
            processed_count += 1
            progress = 25 + int((processed_count / total_to_process) * 65)
            self.update_progress(input_path, progress)
        
        doc.close()
        
        # 确保至少有一张幻灯片
        if len(prs.slides) == 0:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        
        self.update_progress(input_path, 95)
        
        # 保存 PPT
        prs.save(output_path)
        
        self.update_progress(input_path, 100)
        
        return {
            'success': True,
            'output_path': output_path,
            'size': self.get_output_size(output_path),
            'page_count': total_pages,
            'method': 'pymupdf_image'
        }
    
    def _convert_with_pdfplumber(self, input_path: str, output_path: str, 
                                 prs, pages=None) -> Dict[str, Any]:
        """使用 pdfplumber 转换（降级方案 - 文本模式）"""
        import pdfplumber
        from pptx.util import Inches, Pt
        
        with pdfplumber.open(input_path) as pdf:
            if len(pdf.pages) == 0:
                raise Exception("PDF 文件为空，没有可转换的页面")
            
            total_pages = len(pdf.pages)
            
            # 确定要处理的页面
            if pages is None:
                pages_to_process = range(total_pages)
            else:
                pages_to_process = pages
            
            self.update_progress(input_path, 30)
            
            processed_count = 0
            total_to_process = len(pages_to_process)
            
            for page_idx in pages_to_process:
                if page_idx < 0 or page_idx >= total_pages:
                    continue
                
                page = pdf.pages[page_idx]
                
                # 创建新幻灯片（空白布局）
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                # 提取文本
                text = page.extract_text()
                
                if text and text.strip():
                    # 添加文本框
                    left = Inches(0.3)
                    top = Inches(0.3)
                    width = Inches(prs.slide_width.inches - 0.6)
                    height = Inches(prs.slide_height.inches - 0.6)
                    
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    
                    # 根据文本长度动态调整字体大小
                    text_length = len(text)
                    if text_length > 2000:
                        font_size = Pt(8)
                    elif text_length > 1000:
                        font_size = Pt(9)
                    elif text_length > 500:
                        font_size = Pt(10)
                    else:
                        font_size = Pt(11)
                    
                    # 添加文本段落
                    p = text_frame.paragraphs[0]
                    p.text = text.strip()
                    p.font.size = font_size
                
                # 更新进度 (30% ~ 90%)
                processed_count += 1
                progress = 30 + int((processed_count / total_to_process) * 60)
                self.update_progress(input_path, progress)
        
        # 确保至少有一张幻灯片
        if len(prs.slides) == 0:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        
        self.update_progress(input_path, 95)
        
        # 保存 PPT
        prs.save(output_path)
        
        self.update_progress(input_path, 100)
        
        return {
            'success': True,
            'output_path': output_path,
            'size': self.get_output_size(output_path),
            'page_count': total_pages,
            'method': 'pdfplumber_text'
        }
