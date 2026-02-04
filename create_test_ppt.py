"""
创建一个简单的测试PPT文件
"""

import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片大小为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 幻灯片1：标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "This is a test PPT for video conversion"
    
    # 幻灯片2：内容页
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    title.text = "Slide 2"
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "This is the second slide"
    
    # 幻灯片3：内容页
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    title.text = "Slide 3"
    content = slide3.placeholders[1]
    tf = content.text_frame
    tf.text = "This is the third slide"
    
    # 保存文件
    output_file = "test_presentation.pptx"
    prs.save(output_file)
    
    print(f"✓ 成功创建测试PPT: {output_file}")
    print(f"  - 共3页幻灯片")
    print(f"  - 文件大小: {os.path.getsize(output_file)} bytes")
    print(f"\n请使用这个文件测试转换功能！")
    
except ImportError:
    print("错误：需要安装 python-pptx 库")
    print("请运行: pip install python-pptx")
except Exception as e:
    print(f"创建PPT失败: {e}")
    import traceback
    traceback.print_exc()
