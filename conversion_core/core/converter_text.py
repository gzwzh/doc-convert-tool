#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文本转换器
"""

import os
from pathlib import Path
from conversion_core.core.base import Converter


class TextConverter(Converter):
    """文本转换器"""
    
    def convert(self, file_path):
        """转换文件为文本"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.pdf':
                return self._convert_pdf_to_text(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._convert_excel_to_text(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._convert_ppt_to_text(file_path)
            elif ext in ['.docx', '.doc']:
                return self._convert_word_to_text(file_path)
            else:
                return {'success': False, 'error': f'不支持的文件格式: {ext}'}
        
        except Exception as e:
            return {'success': False, 'error': f'转换失败: {str(e)}'}
    
    def _convert_pdf_to_text(self, file_path):
        """PDF转文本 - 改进版本，更好地处理数学公式和特殊字符"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        # 推送开始进度
        self.update_progress(file_path, 5)
        
        # 尝试使用PyMuPDF（fitz）提取文本，它对数学公式和特殊字符支持更好
        try:
            import fitz  # PyMuPDF
            return self._convert_pdf_to_text_with_pymupdf(file_path)
        except ImportError:
            pass  # 如果PyMuPDF不可用，使用pdfplumber
        
        try:
            import pdfplumber
        except ImportError as e:
            return {'success': False, 'error': f'pdfplumber库导入失败: {str(e)}'}
        
        self.update_progress(file_path, 10)
        
        try:
            text_content = []
            
            with pdfplumber.open(file_path) as pdf:
                if len(pdf.pages) == 0:
                    return {'success': False, 'error': 'PDF文件为空，没有可转换的页面'}
                
                total_pages = len(pdf.pages)
                self.update_progress(file_path, 20)
                
                for page_idx, page in enumerate(pdf.pages, 1):
                    # 更新进度 (20% ~ 80%)
                    progress = 20 + int((page_idx / total_pages) * 60)
                    self.update_progress(file_path, progress)
                    
                    # 提取文本 - 使用更好的参数
                    text = page.extract_text(x_tolerance=2, y_tolerance=2)
                    if text:
                        # 清理文本，处理特殊字符
                        cleaned_text = self._clean_extracted_text(text)
                        text_content.append(f"=== 第 {page_idx} 页 ===\n")
                        text_content.append(cleaned_text)
                        text_content.append("\n\n")
                    
                    # 提取表格
                    tables = page.extract_tables()
                    for table_idx, table in enumerate(tables, 1):
                        if table and len(table) > 0:
                            text_content.append(f"--- 表格 {table_idx} ---\n")
                            for row in table:
                                if row:
                                    row_text = '\t'.join([str(cell) if cell else '' for cell in row])
                                    text_content.append(row_text + '\n')
                            text_content.append('\n')
            
            # 保存文本文件
            self.update_progress(file_path, 85)
            output_file = self.get_output_filename(file_path, 'txt')
            
            self.update_progress(file_path, 90)
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(''.join(text_content))
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'PDF转文本失败: {error_msg}'}
    
    def _convert_pdf_to_text_with_pymupdf(self, file_path):
        """使用PyMuPDF提取文本 - 对数学公式和特殊字符支持更好"""
        import fitz
        
        self.update_progress(file_path, 10)
        
        try:
            text_content = []
            doc = fitz.open(file_path)
            
            if doc.page_count == 0:
                doc.close()
                return {'success': False, 'error': 'PDF文件为空，没有可转换的页面'}
            
            total_pages = doc.page_count
            self.update_progress(file_path, 20)
            
            for page_idx in range(total_pages):
                # 更新进度 (20% ~ 80%)
                progress = 20 + int(((page_idx + 1) / total_pages) * 60)
                self.update_progress(file_path, progress)
                
                page = doc[page_idx]
                
                # 使用"text"模式提取，保留更多格式信息
                # 也可以尝试"dict"模式获取更详细的结构
                text = page.get_text("text", sort=True)
                
                if text and text.strip():
                    # 清理文本
                    cleaned_text = self._clean_extracted_text(text)
                    text_content.append(f"=== 第 {page_idx + 1} 页 ===\n")
                    text_content.append(cleaned_text)
                    text_content.append("\n\n")
            
            doc.close()
            
            # 保存文本文件
            self.update_progress(file_path, 85)
            output_file = self.get_output_filename(file_path, 'txt')
            
            self.update_progress(file_path, 90)
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(''.join(text_content))
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'PDF转文本失败(PyMuPDF): {error_msg}'}
    
    def _clean_extracted_text(self, text):
        """清理提取的文本，处理常见的乱码和特殊字符"""
        if not text:
            return ""
        
        # 常见的数学符号映射（PDF中可能使用特殊编码）
        math_replacements = {
            '\uf0b7': '•',      # 项目符号
            '\uf0d7': '×',      # 乘号
            '\uf0b8': '÷',      # 除号
            '\uf02b': '+',      # 加号
            '\uf02d': '-',      # 减号
            '\uf03d': '=',      # 等号
            '\uf0b1': '±',      # 正负号
            '\uf0a3': '≤',      # 小于等于
            '\uf0b3': '≥',      # 大于等于
            '\uf0b9': '≠',      # 不等于
            '\uf0bb': '∞',      # 无穷
            '\uf070': 'π',      # 圆周率
            '\uf061': 'α',      # alpha
            '\uf062': 'β',      # beta
            '\uf067': 'γ',      # gamma
            '\uf064': 'δ',      # delta
            '\uf065': 'ε',      # epsilon
            '\uf071': 'θ',      # theta
            '\uf06c': 'λ',      # lambda
            '\uf06d': 'μ',      # mu
            '\uf073': 'σ',      # sigma
            '\uf077': 'ω',      # omega
            '\uf044': 'Δ',      # Delta
            '\uf053': 'Σ',      # Sigma
            '\uf050': 'Π',      # Pi
            '\uf057': 'Ω',      # Omega
            '\uf0d6': '√',      # 根号
            '\uf0b6': '∂',      # 偏导
            '\uf0f2': '∫',      # 积分
            '\uf0e5': '∑',      # 求和
            '\uf0d5': '∏',      # 连乘
            '\uf0ce': '∈',      # 属于
            '\uf0cf': '∉',      # 不属于
            '\uf0cc': '⊂',      # 子集
            '\uf0c8': '∩',      # 交集
            '\uf0c7': '∪',      # 并集
            '\uf0c6': '∅',      # 空集
            '\uf0ae': '→',      # 右箭头
            '\uf0ac': '←',      # 左箭头
            '\uf0ab': '↑',      # 上箭头
            '\uf0ad': '↓',      # 下箭头
            '\uf0db': '⇒',      # 双线右箭头
            '\uf0dc': '⇔',      # 双线双向箭头
        }
        
        # 应用替换
        for old, new in math_replacements.items():
            text = text.replace(old, new)
        
        # 移除其他私有区域字符（U+E000 - U+F8FF），但保留常见Unicode字符
        import re
        # 只移除明显的乱码字符，保留可能有意义的特殊字符
        text = re.sub(r'[\ue000-\ue0ff]', '', text)  # 移除私有区域低端
        
        # 清理多余的空白行
        lines = text.split('\n')
        cleaned_lines = []
        prev_empty = False
        for line in lines:
            is_empty = not line.strip()
            if is_empty and prev_empty:
                continue  # 跳过连续的空行
            cleaned_lines.append(line)
            prev_empty = is_empty
        
        return '\n'.join(cleaned_lines)
    
    def _convert_excel_to_text(self, file_path):
        """Excel转文本"""
        self.update_progress(file_path, 10)
        
        try:
            from openpyxl import load_workbook
        except ImportError as e:
            return {'success': False, 'error': f'openpyxl库导入失败: {str(e)}'}
        
        try:
            self.update_progress(file_path, 20)
            wb = load_workbook(file_path, data_only=True)
            text_content = []
            
            total_sheets = len(wb.sheetnames)
            self.update_progress(file_path, 30)
            
            for sheet_idx, sheet_name in enumerate(wb.sheetnames, 1):
                ws = wb[sheet_name]
                text_content.append(f"=== 工作表: {sheet_name} ===\n\n")
                
                for row in ws.iter_rows(values_only=True):
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    text_content.append(row_text + '\n')
                
                text_content.append('\n\n')
                
                # 更新进度 (30% ~ 80%)
                progress = 30 + int((sheet_idx / total_sheets) * 50)
                self.update_progress(file_path, progress)
            
            self.update_progress(file_path, 85)
            output_file = self.get_output_filename(file_path, 'txt')
            
            self.update_progress(file_path, 90)
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(''.join(text_content))
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            return {'success': False, 'error': f'Excel转文本失败: {str(e)}'}
    
    def _convert_ppt_to_text(self, file_path):
        """PPT转文本"""
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.ppt':
            return {'success': False, 'error': 'PPT格式（旧格式）不支持'}
        
        self.update_progress(file_path, 10)
        
        try:
            from pptx import Presentation
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        try:
            self.update_progress(file_path, 20)
            prs = Presentation(file_path)
            text_content = []
            
            total_slides = len(prs.slides)
            self.update_progress(file_path, 30)
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                text_content.append(f"=== 幻灯片 {slide_idx} ===\n\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text + '\n')
                
                text_content.append('\n\n')
                
                # 更新进度 (30% ~ 80%)
                progress = 30 + int((slide_idx / total_slides) * 50)
                self.update_progress(file_path, progress)
            
            self.update_progress(file_path, 85)
            output_file = self.get_output_filename(file_path, 'txt')
            
            self.update_progress(file_path, 90)
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(''.join(text_content))
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            return {'success': False, 'error': f'PPT转文本失败: {str(e)}'}
    
    def _convert_word_to_text(self, file_path):
        """Word转文本"""
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.doc':
            return {'success': False, 'error': 'DOC格式（旧格式）不支持'}
        
        self.update_progress(file_path, 10)
        
        try:
            from docx import Document
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        try:
            self.update_progress(file_path, 20)
            doc = Document(file_path)
            text_content = []
            
            total_elements = len(doc.paragraphs) + len(doc.tables)
            self.update_progress(file_path, 30)
            
            processed = 0
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text + '\n')
                processed += 1
                
                # 更新进度 (30% ~ 70%)
                if total_elements > 0:
                    progress = 30 + int((processed / total_elements) * 40)
                    self.update_progress(file_path, progress)
            
            # 处理表格
            for table in doc.tables:
                text_content.append('\n--- 表格 ---\n')
                for row in table.rows:
                    row_text = '\t'.join([cell.text for cell in row.cells])
                    text_content.append(row_text + '\n')
                text_content.append('\n')
                processed += 1
                
                # 更新进度 (30% ~ 70%)
                if total_elements > 0:
                    progress = 30 + int((processed / total_elements) * 40)
                    self.update_progress(file_path, progress)
            
            self.update_progress(file_path, 85)
            output_file = self.get_output_filename(file_path, 'txt')
            
            self.update_progress(file_path, 90)
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(''.join(text_content))
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            return {'success': False, 'error': f'Word转文本失败: {str(e)}'}