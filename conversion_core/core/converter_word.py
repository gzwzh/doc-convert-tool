#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Word转换器
"""

import os
from pathlib import Path
import re
from conversion_core.core.base import Converter


class WordConverter(Converter):
    """Word转换器"""
    
    def _calculate_col_widths_pymupdf(self, table) -> list:
        """从PyMuPDF table计算列宽比例"""
        try:
            cells = table.cells
            if not cells:
                return []
            
            # 收集所有x坐标
            x_coords = set()
            for cell in cells:
                x_coords.add(round(cell[0], 1))  # x0
                x_coords.add(round(cell[2], 1))  # x1
            
            x_coords = sorted(x_coords)
            if len(x_coords) < 2:
                return []
            
            # 计算列宽
            total_width = x_coords[-1] - x_coords[0]
            if total_width <= 0:
                return []
            
            col_widths = []
            for i in range(len(x_coords) - 1):
                width = (x_coords[i + 1] - x_coords[i]) / total_width
                # 限制最小列宽为5%，最大为80%
                width = max(0.05, min(0.8, width))
                col_widths.append(width)
            
            return col_widths
        except Exception as e:
            print(f"[列宽计算错误-PyMuPDF] {e}")
            return []
    
    def _calculate_col_widths_pdfplumber(self, table_data: list, page) -> list:
        """从pdfplumber表格数据计算列宽比例（基于内容长度估算）"""
        try:
            if not table_data or len(table_data) == 0:
                return []
            
            # 获取最大列数
            max_cols = max(len(row) for row in table_data if row)
            if max_cols == 0:
                return []
            
            # 基于内容长度估算列宽
            col_max_lengths = [0] * max_cols
            for row in table_data:
                if row:
                    for col_idx, cell in enumerate(row):
                        if col_idx < max_cols and cell:
                            cell_len = len(str(cell).strip())
                            col_max_lengths[col_idx] = max(col_max_lengths[col_idx], cell_len)
            
            total_length = sum(col_max_lengths) or 1
            col_widths = []
            for length in col_max_lengths:
                width = max(length, 1) / total_length
                # 限制最小列宽为5%，最大为80%
                width = max(0.05, min(0.8, width))
                col_widths.append(width)
            
            # 归一化
            total = sum(col_widths)
            if total > 0:
                col_widths = [w / total for w in col_widths]
            
            return col_widths
        except Exception as e:
            print(f"[列宽计算错误-pdfplumber] {e}")
            return []
    
    def _apply_col_widths(self, word_table, col_widths: list, page_width_inches: float = 6.5):
        """应用列宽到Word表格"""
        try:
            from docx.shared import Inches
            
            if not col_widths or len(col_widths) == 0:
                return
            
            for col_idx, width_ratio in enumerate(col_widths):
                if col_idx < len(word_table.columns):
                    col_width = Inches(page_width_inches * width_ratio)
                    for cell in word_table.columns[col_idx].cells:
                        cell.width = col_width
        except Exception as e:
            print(f"[应用列宽错误] {e}")
    
    def _get_table_hash(self, table_data: list) -> str:
        """计算表格内容hash用于去重"""
        import hashlib
        content = ""
        for row in table_data:
            if row:
                content += "|".join(str(cell).strip().lower() for cell in row if cell) + "\n"
        return hashlib.md5(content.encode('utf-8')).hexdigest()
    
    def _deduplicate_tables(self, tables: list) -> list:
        """表格去重"""
        if not tables:
            return []
        
        seen_hashes = set()
        unique_tables = []
        
        for table in tables:
            if not table:
                continue
            table_hash = self._get_table_hash(table)
            if table_hash not in seen_hashes:
                seen_hashes.add(table_hash)
                unique_tables.append(table)
        
        return unique_tables
    
    def _process_image_alpha(self, image_bytes: bytes) -> bytes:
        """处理图片alpha通道，将透明背景合成为白色背景
        
        解决印章等带透明背景的图片在Word中显示为黑色背景的问题
        """
        try:
            from PIL import Image
            from io import BytesIO
            
            img = Image.open(BytesIO(image_bytes))
            
            # 检查是否有alpha通道
            if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                # 创建白色背景
                background = Image.new('RGB', img.size, (255, 255, 255))
                
                # 转换为RGBA以获取alpha通道
                if img.mode != 'RGBA':
                    img = img.convert('RGBA')
                
                # 使用alpha通道作为mask合成
                background.paste(img, mask=img.split()[3])
                img = background
            elif img.mode == 'CMYK':
                # CMYK转RGB
                img = img.convert('RGB')
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            # 输出为PNG
            output = BytesIO()
            img.save(output, format='PNG')
            return output.getvalue()
            
        except Exception as e:
            print(f"[图片Alpha处理] 处理失败，返回原图: {e}")
            return image_bytes
    
    def _extract_tables_for_word(self, page):
        """为Word转换提取表格（多策略）"""
        try:
            tables = []
            
            # 策略1: 默认设置
            default_tables = page.extract_tables()
            if default_tables:
                tables.extend(default_tables)
            
            # 策略2: 严格线条检测（有边框表格）
            strict_settings = {
                "vertical_strategy": "lines",
                "horizontal_strategy": "lines",
                "snap_tolerance": 3,
                "join_tolerance": 3,
                "edge_min_length": 5
            }
            strict_tables = page.extract_tables(strict_settings)
            if strict_tables:
                for table in strict_tables:
                    if table not in tables:
                        tables.append(table)
            
            # 策略3: 文本对齐检测（无边框表格）
            text_settings = {
                "vertical_strategy": "text",
                "horizontal_strategy": "text",
                "snap_tolerance": 5,
                "join_tolerance": 5,
                "text_tolerance": 3,
                "intersection_tolerance": 5
            }
            try:
                text_tables = page.extract_tables(text_settings)
                if text_tables:
                    for table in text_tables:
                        if table and table not in tables and self._is_valid_table(table):
                            tables.append(table)
            except Exception:
                pass
            
            # 策略4: 严格线条+宽松容差
            lines_strict_settings = {
                "vertical_strategy": "lines_strict",
                "horizontal_strategy": "lines_strict",
                "snap_tolerance": 5,
                "join_tolerance": 5,
                "edge_min_length": 3
            }
            try:
                lines_strict_tables = page.extract_tables(lines_strict_settings)
                if lines_strict_tables:
                    for table in lines_strict_tables:
                        if table and table not in tables:
                            tables.append(table)
            except Exception:
                pass
            
            # 清理表格数据
            cleaned_tables = []
            for table in tables:
                cleaned_table = self._clean_word_table(table)
                if cleaned_table:
                    cleaned_tables.append(cleaned_table)
            
            # 去重
            cleaned_tables = self._deduplicate_tables(cleaned_tables)
            
            return cleaned_tables
            
        except Exception as e:
            print(f"[Word表格提取错误] {e}")
            return []
    
    def _is_valid_table(self, table) -> bool:
        """验证表格是否有效"""
        if not table or len(table) < 2:
            return False
        
        non_empty_cells = 0
        total_cells = 0
        
        for row in table:
            if row:
                for cell in row:
                    total_cells += 1
                    if cell and str(cell).strip():
                        non_empty_cells += 1
        
        # 至少30%的单元格有内容
        return total_cells > 0 and (non_empty_cells / total_cells) >= 0.3
    
    def _handle_doc_format_conversion(self, docx_file: str) -> dict:
        """处理DOC格式转换"""
        doc_file = docx_file.replace('.docx', '.doc')
        doc_converted = self._convert_docx_to_doc(docx_file, doc_file)
        
        if doc_converted:
            # 删除临时docx文件
            try:
                os.remove(docx_file)
            except:
                pass
            return {'success': True, 'output_file': doc_file}
        else:
            # 转换失败，返回docx文件
            return {
                'success': True, 
                'output_file': docx_file,
                'warning': 'DOC格式转换失败，已保存为DOCX格式'
            }
    
    def _convert_pdf_to_word_with_pymupdf(self, file_path: str, output_format: str) -> dict:
        """使用PyMuPDF转换PDF为Word（支持图片提取）"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            # PyMuPDF不可用，回退到pdfplumber
            return self._convert_pdf_to_word_with_pdfplumber(file_path, output_format)
        
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            from docx.oxml.ns import qn
            from io import BytesIO
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        try:
            doc = Document()
            
            # 设置文档默认字体
            doc.styles['Normal'].font.name = 'SimSun'
            doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')
            doc.styles['Normal'].font.size = Pt(10.5)
            
            self.update_progress(file_path, 20)
            
            pdf_doc = fitz.open(file_path)
            total_pages = len(pdf_doc)
            
            if total_pages == 0:
                pdf_doc.close()
                return {'success': False, 'error': 'PDF文件为空'}
            
            self.update_progress(file_path, 25)
            
            for page_num in range(total_pages):
                page = pdf_doc[page_num]
                
                # 提取文本块（保留位置信息）
                text_dict = page.get_text("dict")
                
                for block in text_dict["blocks"]:
                    if "lines" in block:  # 文本块
                        text_content = ""
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text_content += span["text"]
                            text_content += "\n"
                        
                        if text_content.strip():
                            para = doc.add_paragraph(text_content.strip())
                            para.style.font.size = Pt(10.5)
                    
                    elif "image" in block or block.get("type") == 1:  # 图片块
                        # 尝试提取图片
                        try:
                            img_rect = fitz.Rect(block["bbox"])
                            # 从页面裁剪图片区域
                            mat = fitz.Matrix(2.0, 2.0)  # 2倍缩放提高质量
                            pix = page.get_pixmap(matrix=mat, clip=img_rect, alpha=True)
                            
                            # 转换为PNG格式并处理alpha通道
                            img_data = pix.tobytes("png")
                            img_data = self._process_image_alpha(img_data)
                            img_stream = BytesIO(img_data)
                            
                            # 计算图片尺寸（最大宽度6英寸）
                            img_width = (img_rect.x1 - img_rect.x0) / 72  # 转换为英寸
                            if img_width > 6:
                                img_width = 6
                            
                            # 添加图片到文档
                            para = doc.add_paragraph()
                            run = para.add_run()
                            run.add_picture(img_stream, width=Inches(img_width))
                            
                            pix = None
                        except Exception as e:
                            print(f"[PDF转Word] 提取图片块失败: {e}")
                
                # 提取页面中的独立图片
                image_list = page.get_images(full=True)
                print(f"[PDF转Word] 第{page_num + 1}页发现 {len(image_list)} 张图片")
                
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        
                        # 提取图片数据
                        base_image = pdf_doc.extract_image(xref)
                        if base_image:
                            image_bytes = base_image["image"]
                            image_ext = base_image["ext"]
                            
                            # 处理alpha通道，解决印章背景变黑问题
                            image_bytes = self._process_image_alpha(image_bytes)
                            img_stream = BytesIO(image_bytes)
                            
                            # 添加图片到文档
                            para = doc.add_paragraph()
                            run = para.add_run()
                            
                            # 尝试获取图片尺寸
                            try:
                                from PIL import Image
                                pil_img = Image.open(BytesIO(image_bytes))
                                width, height = pil_img.size
                                # 计算适当的宽度（最大6英寸）
                                img_width = min(width / 96, 6)  # 假设96 DPI
                                run.add_picture(img_stream, width=Inches(img_width))
                            except:
                                run.add_picture(img_stream, width=Inches(4))
                            
                            print(f"[PDF转Word] 成功添加图片 {img_index + 1}")
                    except Exception as e:
                        print(f"[PDF转Word] 提取图片 {img_index + 1} 失败: {e}")
                
                # 提取表格（带列宽保留）
                try:
                    tables = page.find_tables()
                    for table in tables:
                        table_data = table.extract()
                        if table_data and len(table_data) > 0:
                            max_cols = max(len(row) for row in table_data if row)
                            if max_cols > 0:
                                word_table = doc.add_table(rows=len(table_data), cols=max_cols)
                                word_table.style = 'Table Grid'
                                
                                for row_idx, row in enumerate(table_data):
                                    if row and row_idx < len(word_table.rows):
                                        for col_idx, cell in enumerate(row):
                                            if col_idx < max_cols:
                                                cell_text = str(cell).strip() if cell else ""
                                                word_table.rows[row_idx].cells[col_idx].text = cell_text
                                
                                # 应用列宽
                                col_widths = self._calculate_col_widths_pymupdf(table)
                                if col_widths:
                                    self._apply_col_widths(word_table, col_widths)
                                
                                doc.add_paragraph()
                except Exception as e:
                    print(f"[PDF转Word] 提取表格失败: {e}")
                
                # 更新进度
                progress = 25 + int(((page_num + 1) / total_pages) * 65)
                self.update_progress(file_path, progress)
                
                # 添加分页符
                if page_num < total_pages - 1:
                    doc.add_page_break()
            
            pdf_doc.close()
            
            self.update_progress(file_path, 95)
            output_file = self.get_output_filename(file_path, output_format)
            
            # 处理DOC格式
            if output_format == 'doc':
                temp_docx = output_file.replace('.doc', '.docx')
                doc.save(temp_docx)
                return self._handle_doc_format_conversion(temp_docx)
            else:
                doc.save(output_file)
                return {'success': True, 'output_file': output_file}
                
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'PyMuPDF转换失败: {error_msg}'}
    
    def _convert_pdf_to_word_with_pdfplumber(self, file_path: str, output_format: str) -> dict:
        """使用pdfplumber转换PDF为Word（备用方案）"""
        try:
            import pdfplumber
        except ImportError as e:
            return {'success': False, 'error': f'pdfplumber库导入失败: {str(e)}'}
        
        try:
            from docx import Document
            from docx.shared import Pt
            from docx.oxml.ns import qn
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        try:
            doc = Document()
            
            doc.styles['Normal'].font.name = 'SimSun'
            doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')
            doc.styles['Normal'].font.size = Pt(10.5)
            
            self.update_progress(file_path, 20)
            
            with pdfplumber.open(file_path) as pdf:
                if len(pdf.pages) == 0:
                    return {'success': False, 'error': 'PDF文件为空'}
                
                total_pages = len(pdf.pages)
                self.update_progress(file_path, 30)
                
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        para = doc.add_paragraph(text)
                        para.style.font.size = Pt(10.5)
                    
                    # 提取表格（带列宽保留）
                    tables = self._extract_tables_for_word(page)
                    for table in tables:
                        if table and len(table) > 0:
                            max_cols = max(len(row) for row in table if row)
                            if max_cols > 0:
                                word_table = doc.add_table(rows=len(table), cols=max_cols)
                                word_table.style = 'Table Grid'
                                
                                for row_idx, row in enumerate(table):
                                    if row_idx < len(word_table.rows):
                                        for col_idx, cell in enumerate(row):
                                            if col_idx < max_cols:
                                                cell_text = str(cell).strip() if cell else ""
                                                word_table.rows[row_idx].cells[col_idx].text = cell_text
                                
                                # 应用列宽
                                col_widths = self._calculate_col_widths_pdfplumber(table, page)
                                if col_widths:
                                    self._apply_col_widths(word_table, col_widths)
                                
                                doc.add_paragraph()
                    
                    progress = 30 + int((page_num / total_pages) * 60)
                    self.update_progress(file_path, progress)
                    
                    if page_num < total_pages:
                        doc.add_page_break()
            
            self.update_progress(file_path, 95)
            output_file = self.get_output_filename(file_path, output_format)
            
            if output_format == 'doc':
                temp_docx = output_file.replace('.doc', '.docx')
                doc.save(temp_docx)
                return self._handle_doc_format_conversion(temp_docx)
            else:
                doc.save(output_file)
                return {'success': True, 'output_file': output_file}
                
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'pdfplumber转换失败: {error_msg}'}
    
    def _clean_word_table(self, table):
        """清理Word表格数据"""
        if not table or len(table) < 1:
            return None
        
        cleaned_table = []
        for row in table:
            if row:
                cleaned_row = []
                for cell in row:
                    if cell is not None:
                        # 清理单元格内容
                        cell_str = str(cell).strip()
                        # 移除多余的换行和空白
                        cell_str = ' '.join(cell_str.split())
                        cleaned_row.append(cell_str)
                    else:
                        cleaned_row.append("")
                
                # 保留所有行，包括部分为空的行
                cleaned_table.append(cleaned_row)
        
        # 检查表格是否有效（至少有一些非空内容）
        has_content = False
        for row in cleaned_table:
            if any(cell.strip() for cell in row):
                has_content = True
                break
        
        return cleaned_table if has_content else None
    
    def _convert_docx_to_doc(self, docx_path: str, doc_path: str) -> bool:
        """尝试将DOCX转换为DOC格式"""
        try:
            # 尝试使用Windows COM接口
            import win32com.client
            
            word_app = win32com.client.Dispatch("Word.Application")
            word_app.Visible = False
            
            # 打开DOCX文件
            doc = word_app.Documents.Open(os.path.abspath(docx_path))
            
            # 保存为DOC格式 (wdFormatDocument = 0)
            doc.SaveAs2(os.path.abspath(doc_path), FileFormat=0)
            doc.Close()
            word_app.Quit()
            
            return True
            
        except ImportError:
            # win32com不可用，尝试使用LibreOffice
            return self._convert_with_libreoffice(docx_path, doc_path)
        except Exception as e:
            print(f"[DOC转换失败] COM接口错误: {e}")
            return self._convert_with_libreoffice(docx_path, doc_path)
    
    def _convert_with_libreoffice(self, docx_path: str, doc_path: str) -> bool:
        """使用LibreOffice转换DOCX到DOC"""
        try:
            import subprocess
            
            # 尝试找到LibreOffice
            libreoffice_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                "soffice",  # 如果在PATH中
                os.path.join(os.path.dirname(__file__), "..", "..", "..", "resources", "libreoffice", "App", "libreoffice", "program", "soffice.exe")
            ]
            
            soffice_path = None
            for path in libreoffice_paths:
                if os.path.exists(path) or path == "soffice":
                    soffice_path = path
                    break
            
            if not soffice_path:
                return False
            
            # 构建转换命令
            output_dir = os.path.dirname(doc_path)
            cmd = [
                soffice_path,
                "--headless",
                "--convert-to", "doc",
                "--outdir", output_dir,
                docx_path
            ]
            
            # 执行转换
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            
            # 检查转换结果
            expected_doc = os.path.join(output_dir, os.path.splitext(os.path.basename(docx_path))[0] + ".doc")
            if os.path.exists(expected_doc):
                # 如果生成的文件名不匹配，重命名
                if expected_doc != doc_path:
                    os.rename(expected_doc, doc_path)
                return True
            
            return False
            
        except Exception as e:
            print(f"[DOC转换失败] LibreOffice错误: {e}")
            return False
    
    def convert(self, file_path):
        """转换文件为Word"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.pdf':
                return self._convert_pdf_to_word(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._convert_excel_to_word(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._convert_ppt_to_word(file_path)
            elif ext in ['.docx', '.doc']:
                return self._convert_word_to_word(file_path)
            elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff', '.webp']:
                return self._convert_image_to_word(file_path)
            else:
                return {'success': False, 'error': f'不支持的文件格式: {ext}'}
        
        except Exception as e:
            return {'success': False, 'error': f'转换失败: {str(e)}'}
    
    def _parse_page_range(self, page_input, total_pages):
        """解析页码范围，返回合法且在总页数范围内的页码列表"""
        if not page_input:
            return list(range(1, total_pages + 1))

        normalized = str(page_input).replace('，', ',').replace('－', '-').replace('–', '-').replace('—', '-')
        parts = normalized.replace(' ', '').split(',')
        pages = set()

        for part in parts:
            if not part:
                continue
            if '-' in part:
                start_str, end_str = part.split('-', 1)
                try:
                    start = int(start_str)
                    end = int(end_str)
                except ValueError:
                    continue
                if start < 1 or end < 1 or end < start:
                    continue
                if start > total_pages or end > total_pages:
                    raise ValueError(f"页码范围超出文档页数: 文档共 {total_pages} 页，范围 {start}-{end}")
                for p in range(start, end + 1):
                    pages.add(p)
            else:
                try:
                    num = int(part)
                except ValueError:
                    continue
                if num < 1:
                    continue
                if num > total_pages:
                    raise ValueError(f"页码超出文档页数: 文档共 {total_pages} 页，页码 {num}")
                pages.add(num)

        if not pages:
            raise ValueError("未解析到有效页码")

        return sorted(pages)
    
    def _convert_pdf_to_word(self, file_path):
        """PDF转Word - 支持图片提取"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        output_format = self.config.get('format', 'docx')
        mode = self.config.get('mode', 'preserve')
        # 支持两种参数命名方式
        page_mode = self.config.get('page_mode') or self.config.get('pageMode', 'all')
        page_input = self.config.get('page_input') or self.config.get('pageInput', '')
        
        # 推送开始进度
        self.update_progress(file_path, 5)
        
        # 首先尝试使用pdf2docx（如果可用且模式为preserve）
        # pdf2docx能更好地保留格式和图片
        if mode == 'preserve':
            try:
                from pdf2docx import Converter as Pdf2DocxConverter
                pdf2docx_available = True
            except ImportError:
                pdf2docx_available = False
            
            if pdf2docx_available:
                # 使用pdf2docx进行格式保留转换（包含图片）
                try:
                    output_file = self.get_output_filename(file_path, output_format)
                    
                    self.update_progress(file_path, 15)
                    
                    # 处理页码选择
                    if page_mode != 'all':
                        try:
                            from pypdf import PdfReader
                            with open(file_path, 'rb') as f:
                                reader = PdfReader(f)
                                total_pages = len(reader.pages)
                        except Exception as e:
                            return {'success': False, 'error': f'读取PDF页数失败: {str(e)}'}
                        
                        # 解析页码
                        if page_mode == 'all':
                            pages = None
                        else:
                            try:
                                pages = self._parse_page_range(page_input, total_pages)
                            except ValueError as e:
                                return {'success': False, 'error': str(e)}
                            pages = [p - 1 for p in pages]
                        
                        self.update_progress(file_path, 30)
                        cv = Pdf2DocxConverter(file_path)
                        self.update_progress(file_path, 50)
                        cv.convert(output_file, pages=pages)
                        self.update_progress(file_path, 90)
                        cv.close()
                        
                        # 处理DOC格式转换
                        if output_format == 'doc':
                            return self._handle_doc_format_conversion(output_file)
                        return {'success': True, 'output_file': output_file}
                    else:
                        # 全部页面
                        self.update_progress(file_path, 30)
                        cv = Pdf2DocxConverter(file_path)
                        self.update_progress(file_path, 50)
                        cv.convert(output_file)
                        self.update_progress(file_path, 90)
                        cv.close()
                        
                        # 处理DOC格式转换
                        if output_format == 'doc':
                            return self._handle_doc_format_conversion(output_file)
                        return {'success': True, 'output_file': output_file}
                except FileNotFoundError:
                    return {'success': False, 'error': 'PDF文件不存在或无法访问'}
                except PermissionError as e:
                    # 权限错误，可能是加密PDF
                    print(f"[PDF转Word] pdf2docx权限错误: {e}, 尝试备用方法")
                    pass
                except ValueError as e:
                    # 值错误，可能是PDF结构问题
                    print(f"[PDF转Word] pdf2docx值错误: {e}, 尝试备用方法")
                    pass
                except Exception as e:
                    # pdf2docx转换失败，回退到PyMuPDF方法
                    error_msg = str(e)
                    error_type = type(e).__name__
                    print(f"[PDF转Word] pdf2docx失败({error_type}): {error_msg[:100]}, 尝试PyMuPDF方法")
                    pass
        
        # 尝试使用PyMuPDF（支持图片提取）
        try:
            return self._convert_pdf_to_word_with_pymupdf(file_path, output_format)
        except Exception as e:
            print(f"[PDF转Word] PyMuPDF失败: {str(e)[:100]}, 回退到pdfplumber")
            pass
        
        # 使用pdfplumber进行文本提取（最后的备用方案）
        return self._convert_pdf_to_word_with_pdfplumber(file_path, output_format)
    
    def _convert_excel_to_word(self, file_path):
        """Excel转Word"""
        try:
            from openpyxl import load_workbook
        except ImportError as e:
            return {'success': False, 'error': f'openpyxl库导入失败: {str(e)}'}
        
        try:
            from docx import Document
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        output_format = self.config.get('format', 'docx')
        
        try:
            wb = load_workbook(file_path, data_only=True)
            doc = Document()
            
            # 设置文档默认字体为宋体（SimSun）
            from docx.shared import Pt
            from docx.oxml.ns import qn
            
            # 设置中文字体为宋体，英文字体为Times New Roman
            doc.styles['Normal'].font.name = 'SimSun'
            doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')
            doc.styles['Normal'].font.size = Pt(10.5)
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # 添加工作表标题
                heading = doc.add_heading(sheet_name, level=1)
                # 设置标题字体
                heading.style.font.name = 'SimHei'  # 黑体
                heading.style._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimHei')
                
                # 添加表格
                table = doc.add_table(rows=1, cols=ws.max_column)
                table.style = 'Light Grid Accent 1'
                
                # 添加表头
                header_cells = table.rows[0].cells
                for col_idx, cell in enumerate(list(ws.iter_cols(values_only=True, max_row=1))[0], 0):
                    cell_text = str(cell) if cell is not None else ''
                    header_cells[col_idx].text = cell_text
                    # 设置表头字体为宋体（与数据行统一）
                    para = header_cells[col_idx].paragraphs[0]
                    if not para.runs:
                        para.add_run(cell_text)
                    run = para.runs[0]
                    run.font.name = 'SimSun'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')
                    run.font.bold = True  # 表头加粗以区分
                    run.font.size = Pt(10.5)
                
                # 添加数据行
                for row in ws.iter_rows(min_row=2, values_only=True):
                    row_cells = table.add_row().cells
                    for col_idx, value in enumerate(row):
                        cell_text = str(value) if value is not None else ''
                        row_cells[col_idx].text = cell_text
                        # 设置数据行字体为宋体
                        para = row_cells[col_idx].paragraphs[0]
                        if not para.runs:
                            para.add_run(cell_text)
                        run = para.runs[0]
                        run.font.name = 'SimSun'
                        run._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')
                        run.font.size = Pt(10.5)
                
                # 添加换行
                doc.add_paragraph()
            
            output_file = self.get_output_filename(file_path, output_format)
            
            # 处理DOC格式
            if output_format == 'doc':
                temp_docx = output_file.replace('.doc', '.docx')
                doc.save(temp_docx)
                
                doc_converted = self._convert_docx_to_doc(temp_docx, output_file)
                
                try:
                    os.remove(temp_docx)
                except:
                    pass
                
                if doc_converted:
                    return {'success': True, 'output_file': output_file}
                else:
                    final_output = output_file.replace('.doc', '.docx')
                    doc.save(final_output)
                    return {
                        'success': True, 
                        'output_file': final_output,
                        'warning': 'DOC格式转换失败，已保存为DOCX格式'
                    }
            else:
                doc.save(output_file)
                return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            import traceback
            error_detail = traceback.format_exc()
            return {'success': False, 'error': f'Excel转Word失败: {str(e)}\n详细信息: {error_detail[:300]}'}
    
    def _convert_ppt_to_word(self, file_path):
        """PPT转Word"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        input_ext = Path(file_path).suffix.lower()
        
        # python-pptx只支持.pptx格式，不支持.ppt格式（旧格式）
        if input_ext == '.ppt':
            return {'success': False, 'error': 'PPT格式（旧格式）不支持。python-pptx库仅支持PPTX格式（新格式）。请先将PPT文件转换为PPTX格式。'}
        
        try:
            from pptx import Presentation
            from pptx.exc import PackageNotFoundError
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        try:
            from docx import Document
            from docx.shared import Pt
            from docx.oxml.ns import qn
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        output_format = self.config.get('format', 'docx')
        
        try:
            prs = Presentation(file_path)
            doc = Document()
            
            # 设置文档默认字体为宋体
            doc.styles['Normal'].font.name = 'SimSun'
            doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')
            doc.styles['Normal'].font.size = Pt(12)
            
            # 设置标题样式为宋体
            doc.styles['Heading 1'].font.name = 'SimSun'
            doc.styles['Heading 1']._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')
            doc.styles['Heading 1'].font.size = Pt(16)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # 添加幻灯片标题
                heading = doc.add_heading(f'幻灯片 {slide_num}', level=1)
                
                # 提取文本内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        paragraph = doc.add_paragraph()
                        
                        # 添加文本并直接设置run的字体为宋体
                        run = paragraph.add_run(shape.text)
                        run.font.name = 'SimSun'
                        if hasattr(run._element, 'rPr') and hasattr(run._element.rPr, 'rFonts'):
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')
                        run.font.size = Pt(12)
                
                # 添加换行
                doc.add_paragraph()
            
            output_file = self.get_output_filename(file_path, output_format)
            
            # 处理DOC格式
            if output_format == 'doc':
                temp_docx = output_file.replace('.doc', '.docx')
                doc.save(temp_docx)
                
                doc_converted = self._convert_docx_to_doc(temp_docx, output_file)
                
                try:
                    os.remove(temp_docx)
                except:
                    pass
                
                if doc_converted:
                    return {'success': True, 'output_file': output_file}
                else:
                    final_output = output_file.replace('.doc', '.docx')
                    doc.save(final_output)
                    return {
                        'success': True, 
                        'output_file': final_output,
                        'warning': 'DOC格式转换失败，已保存为DOCX格式'
                    }
            else:
                doc.save(output_file)
                return {'success': True, 'output_file': output_file}
        
        except PackageNotFoundError:
            return {'success': False, 'error': 'PPT文件格式错误或文件已损坏。请检查文件是否为有效的PPTX格式。'}
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'PPT转Word失败: {error_msg}'}
    
    def _convert_word_to_word(self, file_path):
        """Word转Word（格式转换）"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        try:
            from docx import Document
            from docx.opc.exceptions import PackageNotFoundError
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        output_format = self.config.get('format', 'docx')
        input_ext = Path(file_path).suffix.lower()
        
        # 如果输入输出格式相同，直接复制
        if (input_ext == '.docx' and output_format == 'docx') or (input_ext == '.doc' and output_format == 'doc'):
            try:
                output_file = self.get_output_filename(file_path, output_format)
                import shutil
                shutil.copy2(file_path, output_file)
                return {'success': True, 'output_file': output_file}
            except Exception as e:
                return {'success': False, 'error': f'文件复制失败: {str(e)}'}
        
        # 处理DOC格式（旧格式）- python-docx只支持DOCX
        if input_ext == '.doc':
            return {'success': False, 'error': 'DOC格式（旧格式）需要先转换为DOCX。python-docx库仅支持DOCX格式。'}
        
        # 否则进行格式转换（DOCX格式）
        try:
            doc = Document(file_path)
            
            # 验证文档是否可以读取
            if len(doc.paragraphs) == 0 and len(doc.tables) == 0:
                # 文档为空，但继续转换（创建一个空文档）
                pass
            
            output_file = self.get_output_filename(file_path, output_format)
            
            # python-docx只能保存为DOCX，不能保存为DOC
            if output_format == 'doc':
                return {'success': False, 'error': '无法直接保存为DOC格式，请选择DOCX格式'}
            
            doc.save(output_file)
            return {'success': True, 'output_file': output_file}
        
        except PackageNotFoundError:
            return {'success': False, 'error': 'Word文件格式错误或文件已损坏。请检查文件是否为有效的DOCX格式。'}
        except Exception as e:
            error_msg = str(e)
            # 截断过长的错误信息
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'Word格式转换失败: {error_msg}'}
    
    def _convert_image_to_word(self, file_path):
        """图片转Word"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        try:
            from docx import Document
            from docx.shared import Inches
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        try:
            from PIL import Image
        except ImportError as e:
            return {'success': False, 'error': f'Pillow库导入失败: {str(e)}'}
        
        output_format = self.config.get('format', 'docx')
        
        try:
            # 打开图片
            img = Image.open(file_path)
            
            # 创建Word文档
            doc = Document()
            
            # 将图片插入Word文档
            # 计算适合的图片尺寸（最大宽度6.5英寸）
            max_width = Inches(6.5)
            width, height = img.size
            aspect_ratio = height / width
            pic_width = max_width
            pic_height = Inches(aspect_ratio * 6.5)
            
            # 添加图片到文档
            paragraph = doc.add_paragraph()
            run = paragraph.add_run()
            
            # 保存图片到临时文件
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                img.save(tmp_path, 'PNG')
            
            try:
                run.add_picture(tmp_path, width=pic_width)
            finally:
                # 删除临时文件
                try:
                    os.unlink(tmp_path)
                except:
                    pass
            
            output_file = self.get_output_filename(file_path, output_format)
            
            # 处理DOC格式
            if output_format == 'doc':
                temp_docx = output_file.replace('.doc', '.docx')
                doc.save(temp_docx)
                
                doc_converted = self._convert_docx_to_doc(temp_docx, output_file)
                
                try:
                    os.remove(temp_docx)
                except:
                    pass
                
                if doc_converted:
                    return {'success': True, 'output_file': output_file}
                else:
                    final_output = output_file.replace('.doc', '.docx')
                    doc.save(final_output)
                    return {
                        'success': True, 
                        'output_file': final_output,
                        'warning': 'DOC格式转换失败，已保存为DOCX格式'
                    }
            else:
                doc.save(output_file)
                return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'图片转Word失败: {error_msg}'}
