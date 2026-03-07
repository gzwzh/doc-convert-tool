#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
HTML转换器
"""

import os
from pathlib import Path
from conversion_core.core.base import Converter


class HTMLConverter(Converter):
    """HTML转换器"""
    
    def convert(self, file_path):
        """转换文件为HTML"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.pdf':
                return self._convert_pdf_to_html(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._convert_excel_to_html(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._convert_ppt_to_html(file_path)
            elif ext in ['.docx', '.doc']:
                return self._convert_word_to_html(file_path)
            else:
                return {'success': False, 'error': f'不支持的文件格式: {ext}'}
        
        except Exception as e:
            return {'success': False, 'error': f'转换失败: {str(e)}'}
    
    def _convert_pdf_to_html(self, file_path):
        """PDF转HTML - 使用PyMuPDF保留格式和图片"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        # 推送开始进度
        self.update_progress(file_path, 5)
        
        # 优先使用PyMuPDF（支持图片和格式）
        try:
            import fitz  # PyMuPDF
            pymupdf_available = True
            print(f"[PDF转HTML] 使用PyMuPDF，版本: {fitz.version}")
        except ImportError:
            pymupdf_available = False
            print("[PDF转HTML] PyMuPDF不可用，使用pdfplumber（功能受限）")
        
        self.update_progress(file_path, 10)
        
        try:
            if pymupdf_available:
                return self._convert_pdf_to_html_with_pymupdf(file_path)
            else:
                return self._convert_pdf_to_html_with_pdfplumber(file_path)
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'PDF转HTML失败: {error_msg}'}
    
    def _convert_pdf_to_html_with_pymupdf(self, file_path):
        """使用PyMuPDF转换PDF为HTML（保留格式和图片）"""
        import fitz
        import base64
        
        # 检查是否使用页面渲染模式（将整个页面渲染为图片）
        use_page_render = self.config.get('page_render', True)  # 默认使用页面渲染模式
        
        self.update_progress(file_path, 15)
        
        doc = fitz.open(file_path)
        total_pages = len(doc)
        
        if total_pages == 0:
            doc.close()
            return {'success': False, 'error': 'PDF文件为空，没有可转换的页面'}
        
        self.update_progress(file_path, 20)
        
        # 创建输出目录用于存储图片
        output_file = self.get_output_filename(file_path, 'html')
        output_dir = os.path.dirname(output_file)
        images_dir = os.path.join(output_dir, f"{os.path.splitext(os.path.basename(output_file))[0]}_images")
        os.makedirs(images_dir, exist_ok=True)
        
        html_content = []
        html_content.append('''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PDF转换结果</title>
    <style>
        body { 
            font-family: Arial, sans-serif; 
            margin: 0; 
            padding: 20px; 
            background-color: #f5f5f5;
        }
        .page { 
            background: white;
            margin: 20px auto; 
            padding: 40px;
            box-shadow: 0 0 10px rgba(0,0,0,0.1);
            max-width: 800px;
            position: relative;
            min-height: 1000px;
        }
        .page-header { 
            position: absolute;
            top: 10px;
            right: 15px;
            font-size: 12px; 
            color: #666; 
            background: #f0f0f0;
            padding: 5px 10px;
            border-radius: 3px;
        }
        .text-block {
            position: absolute;
            font-family: inherit;
            white-space: pre-wrap;
            word-wrap: break-word;
        }
        .image-block {
            position: absolute;
            border: 1px solid #ddd;
        }
        .table-block {
            position: absolute;
            border-collapse: collapse;
        }
        .table-block td, .table-block th {
            border: 1px solid #ccc;
            padding: 4px;
            font-size: 12px;
        }
        .table-block th {
            background-color: #f2f2f2;
            font-weight: bold;
        }
    </style>
</head>
<body>
''')
        
        self.update_progress(file_path, 25)
        
        # 如果使用页面渲染模式，将整个页面渲染为图片
        if use_page_render:
            return self._convert_pdf_to_html_page_render(doc, file_path, output_file, images_dir, total_pages)
        
        for page_num in range(total_pages):
            page = doc[page_num]
            
            # 获取页面尺寸
            page_rect = page.rect
            page_width = page_rect.width
            page_height = page_rect.height
            
            # 缩放比例（适应网页显示）
            scale_factor = 800 / page_width if page_width > 800 else 1
            scaled_width = page_width * scale_factor
            scaled_height = page_height * scale_factor
            
            html_content.append(f'<div class="page" style="width: {scaled_width}px; height: {scaled_height}px;">')
            html_content.append(f'<div class="page-header">第 {page_num + 1} 页</div>')
            
            # 提取文本块（保留位置和格式）
            text_dict = page.get_text("dict")
            for block in text_dict["blocks"]:
                if "lines" in block:  # 文本块
                    block_rect = fitz.Rect(block["bbox"])
                    
                    # 计算缩放后的位置
                    left = block_rect.x0 * scale_factor
                    top = block_rect.y0 * scale_factor
                    width = (block_rect.x1 - block_rect.x0) * scale_factor
                    height = (block_rect.y1 - block_rect.y0) * scale_factor
                    
                    # 提取文本内容
                    text_content = ""
                    font_size = 12
                    
                    for line in block["lines"]:
                        for span in line["spans"]:
                            text = span["text"]
                            if text.strip():
                                font_size = max(font_size, span["size"] * scale_factor)
                                text_content += text
                        text_content += "\n"
                    
                    if text_content.strip():
                        # 转义HTML特殊字符
                        text_content = text_content.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        
                        html_content.append(f'''<div class="text-block" style="
                            left: {left}px; 
                            top: {top}px; 
                            width: {width}px; 
                            height: {height}px;
                            font-size: {min(font_size, 16)}px;
                        ">{text_content}</div>''')
            
            # 提取图片
            image_list = page.get_images(full=True)
            print(f"[PDF转HTML] 第{page_num + 1}页发现 {len(image_list)} 张图片")
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    
                    # 获取图片在页面中的位置
                    img_rects = page.get_image_rects(xref)
                    
                    if not img_rects:
                        print(f"[PDF转HTML] 图片 {img_index + 1} 没有找到位置信息，尝试使用默认位置")
                        # 如果没有位置信息，跳过这张图片
                        continue
                    
                    img_rect = img_rects[0]  # 取第一个位置
                    
                    # 计算缩放后的位置
                    left = img_rect.x0 * scale_factor
                    top = img_rect.y0 * scale_factor
                    width = (img_rect.x1 - img_rect.x0) * scale_factor
                    height = (img_rect.y1 - img_rect.y0) * scale_factor
                    
                    # 提取图片数据
                    pix = fitz.Pixmap(doc, xref)
                    
                    # 处理CMYK和其他颜色空间
                    if pix.n - pix.alpha > 3:  # CMYK或其他
                        pix = fitz.Pixmap(fitz.csRGB, pix)  # 转换为RGB
                    
                    # 如果有alpha通道但不是RGBA，转换为RGBA
                    if pix.alpha and pix.n == 2:  # 灰度+alpha
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    
                    # 保存图片
                    img_filename = f"page_{page_num + 1}_img_{img_index + 1}.png"
                    img_path = os.path.join(images_dir, img_filename)
                    pix.save(img_path)
                    print(f"[PDF转HTML] 保存图片: {img_path}")
                    
                    # 相对路径
                    img_rel_path = f"{os.path.basename(images_dir)}/{img_filename}"
                    
                    html_content.append(f'''<img class="image-block" src="{img_rel_path}" style="
                        left: {left}px; 
                        top: {top}px; 
                        width: {width}px; 
                        height: {height}px;
                    " alt="图片 {img_index + 1}">''')
                    
                    pix = None
                except Exception as e:
                    print(f"[PDF转HTML] 处理图片 {img_index + 1} 时出错: {e}")
                    import traceback
                    traceback.print_exc()
                    continue
            
            # 提取表格（如果有）
            try:
                tables = page.find_tables()
                for table_index, table in enumerate(tables):
                    table_rect = table.bbox
                    
                    # 计算缩放后的位置
                    left = table_rect.x0 * scale_factor
                    top = table_rect.y0 * scale_factor
                    width = (table_rect.x1 - table_rect.x0) * scale_factor
                    
                    html_content.append(f'''<table class="table-block" style="
                        left: {left}px; 
                        top: {top}px; 
                        width: {width}px;
                    ">''')
                    
                    # 提取表格数据
                    table_data = table.extract()
                    for row_index, row in enumerate(table_data):
                        tag = 'th' if row_index == 0 else 'td'
                        html_content.append('<tr>')
                        for cell in row:
                            cell_text = str(cell) if cell else ''
                            cell_text = cell_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            html_content.append(f'<{tag}>{cell_text}</{tag}>')
                        html_content.append('</tr>')
                    
                    html_content.append('</table>')
            except:
                # 如果表格提取失败，忽略
                pass
            
            html_content.append('</div>')
            
            # 更新进度 (25% ~ 90%)
            progress = 25 + int(((page_num + 1) / total_pages) * 65)
            self.update_progress(file_path, progress)
        
        doc.close()
        
        html_content.append('</body></html>')
        
        # 保存HTML文件
        self.update_progress(file_path, 95)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
        
        self.update_progress(file_path, 100)
        return {'success': True, 'output_file': output_file}
    
    def _convert_pdf_to_html_page_render(self, doc, file_path, output_file, images_dir, total_pages):
        """将PDF每页渲染为图片嵌入HTML（最可靠的方式）"""
        import fitz
        
        html_content = []
        html_content.append('''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PDF转换结果</title>
    <style>
        body { 
            font-family: Arial, sans-serif; 
            margin: 0; 
            padding: 20px; 
            background-color: #f5f5f5;
            text-align: center;
        }
        .page { 
            background: white;
            margin: 20px auto; 
            padding: 10px;
            box-shadow: 0 0 10px rgba(0,0,0,0.1);
            display: inline-block;
        }
        .page-header { 
            font-size: 14px; 
            color: #666; 
            margin-bottom: 10px;
            text-align: center;
        }
        .page-image {
            max-width: 100%;
            height: auto;
            display: block;
        }
    </style>
</head>
<body>
<h1>PDF文档</h1>
''')
        
        for page_num in range(total_pages):
            page = doc[page_num]
            
            # 高质量渲染页面为图片
            # 使用2倍缩放以获得更清晰的图片
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            
            # 保存图片
            img_filename = f"page_{page_num + 1}.png"
            img_path = os.path.join(images_dir, img_filename)
            pix.save(img_path)
            print(f"[PDF转HTML] 渲染并保存页面 {page_num + 1}: {img_path}")
            
            # 相对路径
            img_rel_path = f"{os.path.basename(images_dir)}/{img_filename}"
            
            html_content.append(f'''<div class="page">
    <div class="page-header">第 {page_num + 1} 页 / 共 {total_pages} 页</div>
    <img class="page-image" src="{img_rel_path}" alt="第 {page_num + 1} 页">
</div>''')
            
            pix = None
            
            # 更新进度 (25% ~ 90%)
            progress = 25 + int(((page_num + 1) / total_pages) * 65)
            self.update_progress(file_path, progress)
        
        doc.close()
        
        html_content.append('</body></html>')
        
        # 保存HTML文件
        self.update_progress(file_path, 95)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
        
        self.update_progress(file_path, 100)
        return {'success': True, 'output_file': output_file}
    
    def _convert_pdf_to_html_with_pdfplumber(self, file_path):
        """使用pdfplumber转换PDF为HTML（功能受限）"""
        try:
            import pdfplumber
        except ImportError as e:
            return {'success': False, 'error': f'pdfplumber库导入失败: {str(e)}'}
        
        self.update_progress(file_path, 15)
        
        html_content = []
        html_content.append('''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PDF转换结果</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; line-height: 1.6; }
        .page { margin-bottom: 30px; padding: 20px; border: 1px solid #ddd; }
        .page-header { font-size: 18px; font-weight: bold; color: #333; margin-bottom: 15px; }
        .table { border-collapse: collapse; width: 100%; margin: 15px 0; }
        .table th, .table td { border: 1px solid #ddd; padding: 8px; text-align: left; }
        .table th { background-color: #f2f2f2; }
        .table-header { font-weight: bold; margin: 10px 0 5px 0; }
        pre { white-space: pre-wrap; word-wrap: break-word; }
    </style>
</head>
<body>
''')
        
        with pdfplumber.open(file_path) as pdf:
            if len(pdf.pages) == 0:
                return {'success': False, 'error': 'PDF文件为空，没有可转换的页面'}
            
            total_pages = len(pdf.pages)
            self.update_progress(file_path, 25)
            
            for page_idx, page in enumerate(pdf.pages, 1):
                html_content.append(f'<div class="page">')
                html_content.append(f'<div class="page-header">第 {page_idx} 页</div>')
                
                # 提取文本
                text = page.extract_text()
                if text:
                    # 转义HTML特殊字符
                    text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    html_content.append(f'<pre>{text}</pre>')
                
                # 提取表格
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables, 1):
                    if table and len(table) > 0:
                        html_content.append(f'<div class="table-header">表格 {table_idx}</div>')
                        html_content.append('<table class="table">')
                        
                        for row_idx, row in enumerate(table):
                            if row:
                                tag = 'th' if row_idx == 0 else 'td'
                                html_content.append('<tr>')
                                for cell in row:
                                    cell_text = str(cell) if cell else ''
                                    cell_text = cell_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                    html_content.append(f'<{tag}>{cell_text}</{tag}>')
                                html_content.append('</tr>')
                        
                        html_content.append('</table>')
                
                html_content.append('</div>')
                
                # 更新进度 (25% ~ 90%)
                progress = 25 + int((page_idx / total_pages) * 65)
                self.update_progress(file_path, progress)
        
        html_content.append('</body></html>')
        
        # 保存HTML文件
        self.update_progress(file_path, 95)
        output_file = self.get_output_filename(file_path, 'html')
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
        
        self.update_progress(file_path, 100)
        return {'success': True, 'output_file': output_file}
    
    def _convert_excel_to_html(self, file_path):
        """Excel转HTML"""
        self.update_progress(file_path, 10)
        
        try:
            from openpyxl import load_workbook
        except ImportError as e:
            return {'success': False, 'error': f'openpyxl库导入失败: {str(e)}'}
        
        try:
            self.update_progress(file_path, 20)
            wb = load_workbook(file_path, data_only=True)
            html_content = []
            
            html_content.append('''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Excel转换结果</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; }
        .sheet { margin-bottom: 30px; }
        .sheet-title { font-size: 18px; font-weight: bold; margin-bottom: 10px; }
        table { border-collapse: collapse; width: 100%; }
        th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
        th { background-color: #f2f2f2; }
    </style>
</head>
<body>
''')
            
            total_sheets = len(wb.sheetnames)
            self.update_progress(file_path, 30)
            
            for sheet_idx, sheet_name in enumerate(wb.sheetnames, 1):
                ws = wb[sheet_name]
                html_content.append(f'<div class="sheet">')
                html_content.append(f'<div class="sheet-title">工作表: {sheet_name}</div>')
                html_content.append('<table>')
                
                first_row = True
                for row in ws.iter_rows(values_only=True):
                    tag = 'th' if first_row else 'td'
                    html_content.append('<tr>')
                    for cell in row:
                        cell_text = str(cell) if cell is not None else ''
                        cell_text = cell_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        html_content.append(f'<{tag}>{cell_text}</{tag}>')
                    html_content.append('</tr>')
                    first_row = False
                
                html_content.append('</table>')
                html_content.append('</div>')
                
                # 更新进度 (30% ~ 80%)
                progress = 30 + int((sheet_idx / total_sheets) * 50)
                self.update_progress(file_path, progress)
            
            html_content.append('</body></html>')
            
            self.update_progress(file_path, 85)
            output_file = self.get_output_filename(file_path, 'html')
            
            self.update_progress(file_path, 90)
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write('\n'.join(html_content))
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            return {'success': False, 'error': f'Excel转HTML失败: {str(e)}'}
    
    def _convert_ppt_to_html(self, file_path):
        """PPT转HTML"""
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.ppt':
            return {'success': False, 'error': 'PPT格式（旧格式）不支持'}
        
        self.update_progress(file_path, 10)
        
        try:
            from pptx import Presentation
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        try:
            self.update_progress(file_path, 20)
            prs = Presentation(file_path)
            html_content = []
            
            html_content.append('''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PPT转换结果</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; }
        .slide { margin-bottom: 30px; padding: 20px; border: 1px solid #ddd; }
        .slide-title { font-size: 18px; font-weight: bold; margin-bottom: 15px; }
        .slide-content { line-height: 1.6; }
        pre { white-space: pre-wrap; word-wrap: break-word; }
    </style>
</head>
<body>
''')
            
            total_slides = len(prs.slides)
            self.update_progress(file_path, 30)
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                html_content.append(f'<div class="slide">')
                html_content.append(f'<div class="slide-title">幻灯片 {slide_idx}</div>')
                html_content.append('<div class="slide-content">')
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        html_content.append(f'<pre>{text}</pre>')
                
                html_content.append('</div>')
                html_content.append('</div>')
                
                # 更新进度 (30% ~ 80%)
                progress = 30 + int((slide_idx / total_slides) * 50)
                self.update_progress(file_path, progress)
            
            html_content.append('</body></html>')
            
            output_file = self.get_output_filename(file_path, 'html')
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write('\n'.join(html_content))
            
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            return {'success': False, 'error': f'PPT转HTML失败: {str(e)}'}
    
    def _convert_word_to_html(self, file_path):
        """Word转HTML"""
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.doc':
            return {'success': False, 'error': 'DOC格式（旧格式）不支持'}
        
        self.update_progress(file_path, 10)
        
        try:
            from docx import Document
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        try:
            self.update_progress(file_path, 20)
            doc = Document(file_path)
            html_content = []
            
            html_content.append('''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Word转换结果</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; line-height: 1.6; }
        table { border-collapse: collapse; width: 100%; margin: 15px 0; }
        th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
        th { background-color: #f2f2f2; }
        .table-title { font-weight: bold; margin: 10px 0 5px 0; }
        pre { white-space: pre-wrap; word-wrap: break-word; }
    </style>
</head>
<body>
''')
            
            total_elements = len(doc.paragraphs) + len(doc.tables)
            self.update_progress(file_path, 30)
            
            processed = 0
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text = para.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    html_content.append(f'<p>{text}</p>')
                processed += 1
                
                # 更新进度 (30% ~ 70%)
                if total_elements > 0:
                    progress = 30 + int((processed / total_elements) * 40)
                    self.update_progress(file_path, progress)
            
            # 处理表格
            for table_idx, table in enumerate(doc.tables, 1):
                html_content.append(f'<div class="table-title">表格 {table_idx}</div>')
                html_content.append('<table>')
                
                first_row = True
                for row in table.rows:
                    tag = 'th' if first_row else 'td'
                    html_content.append('<tr>')
                    for cell in row.cells:
                        cell_text = cell.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        html_content.append(f'<{tag}>{cell_text}</{tag}>')
                    html_content.append('</tr>')
                    first_row = False
                
                html_content.append('</table>')
                processed += 1
                
                # 更新进度 (30% ~ 70%)
                if total_elements > 0:
                    progress = 30 + int((processed / total_elements) * 40)
                    self.update_progress(file_path, progress)
            
            html_content.append('</body></html>')
            
            output_file = self.get_output_filename(file_path, 'html')
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write('\n'.join(html_content))
            
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            return {'success': False, 'error': f'Word转HTML失败: {str(e)}'}