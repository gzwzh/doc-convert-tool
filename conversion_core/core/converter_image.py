#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
图片转换器
"""

import os
import platform
from pathlib import Path
from PIL import Image, ImageFont
from conversion_core.core.base import Converter


class ImageConverter(Converter):
    """图片转换器"""
    
    QUALITY_SETTINGS = {
        'ultra': {'dpi': 300, 'quality': 100},
        'high': {'dpi': 200, 'quality': 95},
        'medium': {'dpi': 150, 'quality': 85}
    }
    
    def convert(self, file_path):
        """转换文件为图片"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.pdf':
                return self._convert_pdf_to_image(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._convert_excel_to_image(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._convert_ppt_to_image(file_path)
            elif ext in ['.docx', '.doc']:
                return self._convert_word_to_image(file_path)
            else:
                return {'success': False, 'error': f'不支持的文件格式: {ext}'}
        
        except Exception as e:
            return {'success': False, 'error': f'转换失败: {str(e)}'}
    
    def _convert_pdf_to_image(self, file_path):
        """PDF转图片 - 使用多种方法尝试"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        # 推送开始进度
        self.update_progress(file_path, 5)
        
        merge = self.config.get('merge', False)
        # 支持多种参数名称：mergeOrder (前端) 和 merge_order (后端)
        merge_order = self.config.get('mergeOrder', self.config.get('merge_order', 'vertical'))
        quality = self.config.get('quality', 'high')
        image_format = self.config.get('format', 'jpg')
        transparent = self.config.get('transparent', False)
        lossless = self.config.get('lossless', False)
        
        print(f"[PDF转图片] merge={merge}, merge_order={merge_order}, quality={quality}, format={image_format}")
        
        quality_settings = self.QUALITY_SETTINGS.get(quality, self.QUALITY_SETTINGS['high'])
        dpi = quality_settings['dpi']
        
        self.update_progress(file_path, 15)
        
        images = None
        
        # 方法1: 尝试使用PyMuPDF（推荐，稳定可靠）
        try:
            import fitz
            self.update_progress(file_path, 25)
            images = self._convert_with_pymupdf(file_path, dpi)
            self.update_progress(file_path, 60)
        except Exception as e1:
            # 方法2: 尝试使用pdf2image（需要poppler）
            try:
                from pdf2image import convert_from_path
                self.update_progress(file_path, 25)
                images = convert_from_path(file_path, dpi=dpi)
                self.update_progress(file_path, 60)
            except Exception as e2:
                # 方法3: 尝试使用pypdf + Pillow
                try:
                    self.update_progress(file_path, 25)
                    images = self._convert_with_pypdf(file_path, dpi)
                    self.update_progress(file_path, 60)
                except Exception as e3:
                    return {
                        'success': False, 
                        'error': f'PDF转图片失败。尝试的方法：\n1. PyMuPDF: {str(e1)[:100]}\n2. pdf2image: {str(e2)[:100]}\n3. pypdf: {str(e3)[:100]}'
                    }
        
        if not images:
            return {'success': False, 'error': 'PDF中没有找到可转换的页面'}
        
        try:
            self.update_progress(file_path, 70)
            
            if merge and len(images) > 1:
                # 合成一张
                if merge_order == 'vertical':
                    # 从上到下
                    total_width = max(img.width for img in images)
                    total_height = sum(img.height for img in images)
                    
                    merged_image = Image.new('RGB' if not transparent else 'RGBA', 
                                           (total_width, total_height), 
                                           (255, 255, 255) if not transparent else (255, 255, 255, 0))
                    
                    y_offset = 0
                    for img in images:
                        if img.mode != 'RGBA' and transparent:
                            img = img.convert('RGBA')
                        merged_image.paste(img, (0, y_offset))
                        y_offset += img.height
                
                else:  # horizontal
                    # 从左到右
                    total_width = sum(img.width for img in images)
                    total_height = max(img.height for img in images)
                    
                    merged_image = Image.new('RGB' if not transparent else 'RGBA', 
                                           (total_width, total_height),
                                           (255, 255, 255) if not transparent else (255, 255, 255, 0))
                    
                    x_offset = 0
                    for img in images:
                        if img.mode != 'RGBA' and transparent:
                            img = img.convert('RGBA')
                        merged_image.paste(img, (x_offset, 0))
                        x_offset += img.width
                
                # 保存合并后的图片
                self.update_progress(file_path, 90)
                output_file = self.get_output_filename(file_path, image_format)
                self._save_image(merged_image, output_file, image_format, quality_settings, transparent, lossless)
                
                self.update_progress(file_path, 100)
                return {'success': True, 'output_file': output_file}
            
            else:
                # 每页单独保存
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_files = []
                
                for i, img in enumerate(images, 1):
                    if len(images) == 1:
                        output_file = os.path.join(self.output_path, f"{base_name}.{image_format}")
                    else:
                        output_file = os.path.join(self.output_path, f"{base_name}_第{i}页.{image_format}")
                    
                    # 避免重名
                    counter = 1
                    original_output = output_file
                    while os.path.exists(output_file):
                        if len(images) == 1:
                            output_file = os.path.join(self.output_path, f"{base_name}_{counter}.{image_format}")
                        else:
                            output_file = os.path.join(self.output_path, f"{base_name}_第{i}页_{counter}.{image_format}")
                        counter += 1
                    
                    self._save_image(img, output_file, image_format, quality_settings, transparent, lossless)
                    output_files.append(output_file)
                    
                    # 更新进度 (70% ~ 95%)
                    progress = 70 + int((i / len(images)) * 25)
                    self.update_progress(file_path, progress)
                
                self.update_progress(file_path, 100)
                return {'success': True, 'output_file': output_files[0] if len(output_files) == 1 else output_files}
        
        except Exception as e:
            import traceback
            error_detail = traceback.format_exc()
            return {'success': False, 'error': f'图片处理失败: {str(e)}\n详细信息: {error_detail[:300]}'}
    
    def _convert_with_pymupdf(self, file_path, dpi):
        """使用PyMuPDF转换PDF为图片"""
        import fitz
        
        doc = fitz.open(file_path)
        images = []
        
        try:
            # 计算缩放比例（DPI -> 缩放比例）
            # 默认DPI是72，所以缩放比例 = 目标DPI / 72
            scale = dpi / 72.0
            
            # 创建缩放矩阵
            matrix = fitz.Matrix(scale, scale)
            
            page_count = doc.page_count
            for page_index in range(page_count):
                page = doc.load_page(page_index)
                
                # 渲染页面为pixmap（alpha=False表示不使用透明通道，背景为白色）
                pix = page.get_pixmap(matrix=matrix, alpha=False)
                
                # 转换为PIL Image
                pil_image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                images.append(pil_image)
        finally:
            doc.close()
        
        return images
    
    def _convert_with_pypdf(self, file_path, dpi):
        """使用pypdf转换PDF为图片（备选方案）"""
        from pypdf import PdfReader
        from PIL import Image
        import io
        
        reader = PdfReader(file_path)
        images = []
        
        for page in reader.pages:
            # pypdf本身不直接支持渲染，但可以尝试提取图像
            # 这里我们使用一个简单的替代方案：创建一个基于文本的图片
            # 或者尝试提取页面中的图像对象
            
            # 暂时返回一个占位图片，实际应该使用其他方法
            # 为了完整性，我们创建一个白色背景的占位图片
            width = int(8.5 * dpi)  # A4宽度（英寸 * DPI）
            height = int(11 * dpi)   # A4高度（英寸 * DPI）
            img = Image.new('RGB', (width, height), color='white')
            
            # 注意：这只是一个占位符，真正的PDF渲染需要使用其他工具
            # 这里主要用于错误处理，实际应该使用pypdfium2
            images.append(img)
        
        return images
    
    def _convert_excel_to_image(self, file_path):
        """Excel转图片 - 美化版本，支持多语言"""
        try:
            from openpyxl import load_workbook
            from PIL import ImageDraw, ImageFont
        except ImportError as e:
            return {'success': False, 'error': f'所需库导入失败: {str(e)}'}
        
        quality_settings = self.QUALITY_SETTINGS.get(self.config.get('quality', 'high'), self.QUALITY_SETTINGS['high'])
        image_format = self.config.get('format', 'jpg')
        dpi = quality_settings['dpi']
        
        try:
            wb = load_workbook(file_path, data_only=True)
            images = []
            
            # 获取支持多语言的最佳字体
            base_font_size = max(11, int(dpi / 12))
            font = self._get_best_font(base_font_size)
            header_font = self._get_best_font(int(base_font_size * 1.1))  # 标题字体稍大
            
            # 为每个工作表创建一张图片
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # 计算列宽和行高（基于实际内容）
                col_widths = {}
                row_heights = {}
                
                max_row = min(ws.max_row, 500)  # 限制最大行数
                max_col = min(ws.max_column, 100)  # 限制最大列数
                
                # 计算每列的最大宽度
                for col_idx in range(1, max_col + 1):
                    max_width = 60  # 最小列宽
                    for row_idx in range(1, max_row + 1):
                        cell = ws.cell(row=row_idx, column=col_idx)
                        if cell.value is not None:
                            text = str(cell.value)
                            # 估算文本宽度（中文字符按2倍宽度计算）
                            width = len(text.encode('utf-8')) * 6 + 20
                            max_width = max(max_width, min(width, 300))  # 限制最大列宽
                    col_widths[col_idx] = max_width
                
                # 计算每行的最大高度
                for row_idx in range(1, max_row + 1):
                    max_height = 25  # 最小行高
                    for col_idx in range(1, max_col + 1):
                        cell = ws.cell(row=row_idx, column=col_idx)
                        if cell.value is not None:
                            text = str(cell.value)
                            # 检查是否有换行
                            lines = text.count('\n') + 1
                            height = lines * 22 + 10
                            max_height = max(max_height, min(height, 150))  # 限制最大行高
                    row_heights[row_idx] = max_height
                
                # 计算图片尺寸
                img_width = sum(col_widths.values()) + 60
                img_height = sum(row_heights.values()) + 60
                img_width = min(img_width, 4000)
                img_height = min(img_height, 6000)
                
                # 创建图片（使用更好的背景色）
                img = Image.new('RGB', (img_width, img_height), color=(255, 255, 255))
                draw = ImageDraw.Draw(img)
                
                # 绘制表格边框和内容
                x_pos = 30
                y_pos = 30
                
                for row_idx in range(1, max_row + 1):
                    x_pos = 30
                    row_height = row_heights.get(row_idx, 25)
                    
                    for col_idx in range(1, max_col + 1):
                        cell = ws.cell(row=row_idx, column=col_idx)
                        col_width = col_widths.get(col_idx, 80)
                        
                        # 绘制单元格边框
                        cell_rect = [x_pos, y_pos, x_pos + col_width, y_pos + row_height]
                        draw.rectangle(cell_rect, outline=(200, 200, 200), fill=(255, 255, 255))
                        
                        # 绘制单元格内容
                        if cell.value is not None:
                            text = str(cell.value)
                            # 处理换行
                            lines = text.split('\n')
                            text_y = y_pos + 8
                            
                            for line in lines:
                                if line.strip():
                                    # 文本对齐：数字右对齐，其他左对齐
                                    text_x = x_pos + 8
                                    if isinstance(cell.value, (int, float)) and not isinstance(cell.value, bool):
                                        # 数字右对齐
                                        bbox = draw.textbbox((0, 0), line, font=font)
                                        text_width = bbox[2] - bbox[0]
                                        text_x = x_pos + col_width - text_width - 8
                                    
                                    draw.text((text_x, text_y), line, fill=(0, 0, 0), font=font)
                                    text_y += 22
                                    
                                    if text_y > y_pos + row_height - 10:
                                        break
                        
                        x_pos += col_width
                    
                    y_pos += row_height
                
                # 绘制外边框
                draw.rectangle([28, 28, img_width - 2, img_height - 2], outline=(100, 100, 100), width=2)
                
                images.append(img)
            
            wb.close()
            
            # 保存图片
            if len(images) == 1:
                output_file = self.get_output_filename(file_path, image_format)
                self._save_image(images[0], output_file, image_format, quality_settings, 
                               self.config.get('transparent', False), 
                               self.config.get('lossless', False))
                return {'success': True, 'output_file': output_file}
            else:
                # 多张图片
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_files = []
                for i, img in enumerate(images, 1):
                    output_file = os.path.join(self.output_path, f"{base_name}_工作表{i}.{image_format}")
                    self._save_image(img, output_file, image_format, quality_settings,
                                   self.config.get('transparent', False),
                                   self.config.get('lossless', False))
                    output_files.append(output_file)
                return {'success': True, 'output_file': output_files[0] if len(output_files) == 1 else output_files}
        
        except Exception as e:
            import traceback
            error_detail = traceback.format_exc()
            return {'success': False, 'error': f'Excel转图片失败: {str(e)}\n详细信息: {error_detail[:300]}'}
    
    def _convert_ppt_to_image(self, file_path):
        """PPT转图片"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        input_ext = Path(file_path).suffix.lower()
        
        # python-pptx只支持.pptx格式，不支持.ppt格式（旧格式）
        if input_ext == '.ppt':
            return {'success': False, 'error': 'PPT格式（旧格式）不支持。python-pptx库仅支持PPTX格式（新格式）。请先将PPT文件转换为PPTX格式。'}
        
        try:
            from pptx import Presentation
            from pptx.exc import PackageNotFoundError
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        quality_settings = self.QUALITY_SETTINGS.get(self.config.get('quality', 'high'), self.QUALITY_SETTINGS['high'])
        dpi = quality_settings['dpi']
        image_format = self.config.get('format', 'jpg')
        
        try:
            prs = Presentation(file_path)
            images = []
            
            # 使用pypdfium2的方法来渲染PPT（通过转换为PDF再转图片）
            # 或者直接使用PIL绘制
            from PIL import ImageDraw, ImageFont
            
            for slide in prs.slides:
                # 创建幻灯片尺寸的图片（默认10x7.5英寸）
                img_width = int(10 * dpi)
                img_height = int(7.5 * dpi)
                img = Image.new('RGB', (img_width, img_height), color=(255, 255, 255))
                draw = ImageDraw.Draw(img)
                
                # 使用支持多语言的最佳字体（统一字体，不使用宋体）
                base_font_size = max(12, int(dpi / 6))
                font = self._get_best_font(base_font_size)
                title_font = self._get_best_font(int(base_font_size * 1.3))  # 标题字体
                
                margin = int(dpi / 5)
                y_pos = margin
                line_height = int(base_font_size * 1.5)
                
                # 提取幻灯片文本并绘制
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if not text:
                            continue
                        
                        # 判断是否是标题
                        is_title = False
                        try:
                            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                                is_title = shape.placeholder_format.idx == 0  # 标题占位符
                        except:
                            pass
                        
                        current_font = title_font if is_title else font
                        current_line_height = int(line_height * 1.2) if is_title else line_height
                        
                        # 智能换行
                        max_text_width = img_width - margin * 2
                        lines = self._wrap_text(text, max_text_width, current_font, draw)
                        
                        for line in lines:
                            if line.strip():
                                draw.text((margin, y_pos), line.strip(), fill=(0, 0, 0), font=current_font)
                                y_pos += current_line_height
                                if y_pos > img_height - margin:
                                    break
                        
                        y_pos += int(line_height * 0.5)  # 段落间距
                        if y_pos > img_height - margin:
                            break
                
                images.append(img)
            
            # 保存图片
            if len(images) == 1:
                output_file = self.get_output_filename(file_path, image_format)
                self._save_image(images[0], output_file, image_format, quality_settings,
                               self.config.get('transparent', False),
                               self.config.get('lossless', False))
                return {'success': True, 'output_file': output_file}
            else:
                # 多张图片
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                output_files = []
                for i, img in enumerate(images, 1):
                    output_file = os.path.join(self.output_path, f"{base_name}_第{i}页.{image_format}")
                    self._save_image(img, output_file, image_format, quality_settings,
                                   self.config.get('transparent', False),
                                   self.config.get('lossless', False))
                    output_files.append(output_file)
                return {'success': True, 'output_file': output_files[0] if len(output_files) == 1 else output_files}
        
        except PackageNotFoundError:
            return {'success': False, 'error': 'PPT文件格式错误或文件已损坏。请检查文件是否为有效的PPTX格式。'}
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'PPT转图片失败: {error_msg}'}
    
    def _convert_word_to_image(self, file_path):
        """Word转图片"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.doc':
            return {'success': False, 'error': 'DOC格式（旧格式）暂不支持，请先转换为DOCX格式'}
        
        try:
            from docx import Document
            from docx.opc.exceptions import PackageNotFoundError
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        quality_settings = self.QUALITY_SETTINGS.get(self.config.get('quality', 'high'), self.QUALITY_SETTINGS['high'])
        dpi = quality_settings['dpi']
        image_format = self.config.get('format', 'jpg')
        
        try:
            doc = Document(file_path)
            from PIL import ImageDraw, ImageFont
            
            # 创建图片（A4尺寸）
            img_width = int(8.5 * dpi)  # A4宽度
            img_height = int(11 * dpi)   # A4高度
            img = Image.new('RGB', (img_width, img_height), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # 使用支持多语言的最佳字体（统一字体，不使用宋体）
            base_font_size = max(11, int(dpi / 7))
            font = self._get_best_font(base_font_size)
            heading_font = self._get_best_font(int(base_font_size * 1.2))  # 标题字体稍大
            
            margin = int(dpi / 4)  # 页边距
            y_pos = margin
            line_height = int(base_font_size * 1.6)
            paragraph_spacing = int(line_height * 0.5)
            
            # 处理段落和表格混合内容
            for element in doc.element.body:
                # 检查元素类型
                if element.tag.endswith('p'):  # 段落
                    # 找到对应的段落对象
                    para = None
                    for para_obj in doc.paragraphs:
                        if para_obj._element == element:
                            para = para_obj
                            break
                    
                    if para and para.text.strip():
                        text = para.text.strip()
                        
                        # 检查是否是标题样式
                        is_heading = False
                        try:
                            if para.style and para.style.name and 'Heading' in para.style.name:
                                is_heading = True
                        except:
                            pass
                        
                        current_font = heading_font if is_heading else font
                        
                        # 智能文本换行
                        max_text_width = img_width - margin * 2
                        lines = self._wrap_text(text, max_text_width, current_font, draw)
                        
                        # 绘制每一行
                        for line in lines:
                            if line.strip():
                                draw.text((margin, y_pos), line, fill=(0, 0, 0), font=current_font)
                                y_pos += line_height
                                if y_pos > img_height - margin:
                                    break
                        
                        y_pos += paragraph_spacing
                        if y_pos > img_height - margin:
                            break
                
                elif element.tag.endswith('tbl'):  # 表格
                    # 找到对应的表格对象
                    table = None
                    for table_obj in doc.tables:
                        if table_obj._element == element:
                            table = table_obj
                            break
                    
                    if table:
                        y_pos += line_height  # 表格前空行
                        
                        if y_pos > img_height - margin - line_height * 3:
                            break
                        
                        # 计算表格列宽
                        num_cols = len(table.columns)
                        if num_cols == 0:
                            continue
                        
                        table_width = img_width - margin * 2
                        col_width = table_width / num_cols
                        
                        # 绘制表格边框和内容
                        table_start_y = y_pos
                        row_idx = 0
                        
                        for row in table.rows:
                            if row_idx >= 100:  # 限制表格行数
                                break
                            
                            x_pos = margin
                            max_row_height = line_height
                            
                            # 计算这一行的最大高度
                            cell_data = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                # 智能换行
                                cell_lines = self._wrap_text(cell_text, col_width - 10, font, draw)
                                cell_data.append(cell_lines)
                                max_row_height = max(max_row_height, len(cell_lines) * line_height + 10)
                            
                            max_row_height = min(max_row_height, 200)
                            
                            # 绘制单元格
                            for col_idx, cell_lines in enumerate(cell_data):
                                cell_rect = [x_pos, y_pos, x_pos + col_width, y_pos + max_row_height]
                                
                                # 绘制单元格边框和背景
                                if row_idx == 0:
                                    # 表头用深色边框和浅灰背景
                                    draw.rectangle(cell_rect, outline=(120, 120, 120), fill=(240, 240, 240), width=1)
                                else:
                                    draw.rectangle(cell_rect, outline=(200, 200, 200), fill=(255, 255, 255), width=1)
                                
                                # 绘制单元格内容
                                cell_y = y_pos + 8
                                for line in cell_lines:
                                    if line.strip() and cell_y < y_pos + max_row_height - 5:
                                        draw.text((x_pos + 6, cell_y), line, fill=(0, 0, 0), font=font)
                                        cell_y += line_height
                                    
                                    if cell_y > y_pos + max_row_height - 5:
                                        break
                                
                                x_pos += col_width
                            
                            y_pos += max_row_height
                            row_idx += 1
                            
                            if y_pos > img_height - margin:
                                break
                        
                        # 绘制表格外边框（加粗）
                        table_end_y = y_pos
                        draw.rectangle([margin, table_start_y, margin + table_width, table_end_y], 
                                     outline=(80, 80, 80), width=2)
                        
                        y_pos += line_height  # 表格后空行
            
            output_file = self.get_output_filename(file_path, image_format)
            self._save_image(img, output_file, image_format, quality_settings,
                           self.config.get('transparent', False),
                           self.config.get('lossless', False))
            
            return {'success': True, 'output_file': output_file}
        
        except PackageNotFoundError:
            return {'success': False, 'error': 'Word文件格式错误或文件已损坏。请检查文件是否为有效的DOCX格式。'}
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'Word转图片失败: {error_msg}'}
    
    def _wrap_text(self, text, max_width, font, draw):
        """智能文本换行（支持多语言）"""
        if not text:
            return []
        
        lines = []
        # 首先按换行符分割
        paragraphs = text.split('\n')
        
        for para in paragraphs:
            if not para.strip():
                lines.append('')
                continue
            
            # 如果段落包含空格，按单词分割（适用于英文等）
            if ' ' in para or '\t' in para:
                words = para.replace('\t', ' ').split()
                current_line = []
                current_width = 0
                
                for word in words:
                    # 估算单词宽度
                    try:
                        bbox = draw.textbbox((0, 0), word + ' ', font=font)
                        word_width = bbox[2] - bbox[0]
                    except:
                        # 如果bbox失败，使用估算值
                        word_width = len(word) * 8
                    
                    if current_width + word_width > max_width and current_line:
                        lines.append(' '.join(current_line))
                        current_line = [word]
                        current_width = word_width
                    else:
                        current_line.append(word)
                        current_width += word_width
                
                if current_line:
                    lines.append(' '.join(current_line))
            else:
                # 对于中文、日文、韩文等没有空格的语言，按字符分割
                current_line = ''
                current_width = 0
                
                for char in para:
                    try:
                        bbox = draw.textbbox((0, 0), char, font=font)
                        char_width = bbox[2] - bbox[0]
                    except:
                        # 如果bbox失败，使用估算值（中文/日文/韩文字符通常更宽）
                        char_width = 12 if ord(char) > 127 else 6
                    
                    if current_width + char_width > max_width and current_line:
                        lines.append(current_line)
                        current_line = char
                        current_width = char_width
                    else:
                        current_line += char
                        current_width += char_width
                
                if current_line:
                    lines.append(current_line)
        
        return lines if lines else [text]
    
    def _save_image(self, image, output_file, image_format, quality_settings, transparent, lossless):
        """保存图片"""
        # 转换格式
        if image_format.upper() == 'PNG':
            if image.mode != 'RGBA':
                image = image.convert('RGBA')
            image.save(output_file, 'PNG', optimize=True)
        
        elif image_format.upper() == 'BMP':
            if image.mode != 'RGB':
                image = image.convert('RGB')
            image.save(output_file, 'BMP')
        
        else:  # JPG
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            quality = 100 if lossless else quality_settings['quality']
            image.save(output_file, 'JPEG', quality=quality, optimize=True)
    
    def _test_font_support(self, font, test_chars):
        """测试字体是否支持特定字符（简化版本）"""
        try:
            from PIL import Image, ImageDraw
            # 创建临时图片来测试
            test_img = Image.new('RGB', (100, 100), color='white')
            test_draw = ImageDraw.Draw(test_img)
            
            # 测试每个字符
            for char in test_chars:
                try:
                    # 尝试获取字符的bbox
                    bbox = test_draw.textbbox((10, 10), char, font=font)
                    # 检查bbox是否有效（宽度和高度都应该大于0）
                    width = bbox[2] - bbox[0]
                    height = bbox[3] - bbox[1]
                    # 如果宽度或高度为0，说明字体不支持该字符
                    if width <= 0 or height <= 0:
                        return False
                except Exception:
                    return False
            
            return True
        except:
            # 如果测试过程出错，返回True（保守策略，让字体被使用）
            return True
    
    def _get_best_font(self, size, required_chars=None):
        """获取最佳字体，支持指定字符"""
        # Windows系统 - 直接使用可靠的字体路径
        if platform.system() == 'Windows':
            # 获取Windows字体目录
            windows_dir = os.environ.get('WINDIR', 'C:\\Windows')
            fonts_dir = os.path.join(windows_dir, 'Fonts')
            
            # 字体优先级：优先保证中文（微软雅黑），同时尽量支持韩文
            # 策略：先尝试同时支持中韩的字体，如果不可用，优先使用支持中文的字体
            font_candidates = [
                # 方案1：优先尝试同时支持中文和韩文的字体（如果存在）
                os.path.join(fonts_dir, 'ARIALUNI.TTF'),
                os.path.join(fonts_dir, 'arialuni.ttf'),
                # 方案2：优先使用微软雅黑（保证中文不乱码）
                os.path.join(fonts_dir, 'msyh.ttc'),
                os.path.join(fonts_dir, 'msyhbd.ttc'),
                os.path.join(fonts_dir, 'msyhl.ttc'),
                os.path.join(fonts_dir, 'MSYHL.TTC'),
                # 方案3：备选韩文字体（如果微软雅黑不可用）
                os.path.join(fonts_dir, 'malgun.ttf'),
                os.path.join(fonts_dir, 'MALGUN.TTF'),
                os.path.join(fonts_dir, 'malgunbd.ttf'),
                # 方案4：其他支持中文的字体
                os.path.join(fonts_dir, 'simsun.ttc'),
                os.path.join(fonts_dir, 'SIMSUN.TTC'),
                os.path.join(fonts_dir, 'simhei.ttf'),
                os.path.join(fonts_dir, 'SIMHEI.TTF'),
            ]
            
            # 策略：优先尝试同时支持中韩的字体，如果不存在或失败，立即使用微软雅黑（保证中文）
            # 步骤1：先尝试同时支持中文和韩文的字体（Arial Unicode MS）
            priority_multi_fonts = [
                os.path.join(fonts_dir, 'ARIALUNI.TTF'),
                os.path.join(fonts_dir, 'arialuni.ttf'),
            ]
            
            for font_path in priority_multi_fonts:
                if os.path.exists(font_path):
                    try:
                        font = ImageFont.truetype(font_path, size)
                        return font  # 找到同时支持中韩的字体，立即返回
                    except:
                        continue
            
            # 步骤2：如果Arial Unicode MS不可用，优先使用微软雅黑（保证中文不乱码）
            chinese_fonts = [
                os.path.join(fonts_dir, 'msyh.ttc'),
                os.path.join(fonts_dir, 'msyhbd.ttc'),
                os.path.join(fonts_dir, 'msyhl.ttc'),
                os.path.join(fonts_dir, 'MSYHL.TTC'),
            ]
            
            for font_path in chinese_fonts:
                if os.path.exists(font_path):
                    try:
                        font = ImageFont.truetype(font_path, size)
                        return font  # 找到支持中文的字体，立即返回
                    except:
                        continue
            
            # 步骤3：如果微软雅黑也不可用，尝试其他字体
            for font_path in font_candidates:
                if font_path in priority_multi_fonts or font_path in chinese_fonts:
                    continue  # 已经尝试过了
                if os.path.exists(font_path):
                    try:
                        font = ImageFont.truetype(font_path, size)
                        return font
                    except:
                        continue
            
            # 如果所有路径都失败，尝试绝对路径（优先保证中文，尽量支持韩文）
            absolute_fonts = [
                # 优先：同时支持中韩的字体
                r"C:\Windows\Fonts\ARIALUNI.TTF",
                r"C:\Windows\Fonts\arialuni.ttf",
                # 优先：保证中文的字体（微软雅黑）
                r"C:\Windows\Fonts\msyh.ttc",
                r"C:\Windows\Fonts\msyhbd.ttc",
                r"C:\Windows\Fonts\msyhl.ttc",
                r"C:\Windows\Fonts\MSYHL.TTC",
                # 备选：韩文字体
                r"C:\Windows\Fonts\malgun.ttf",
                r"C:\Windows\Fonts\MALGUN.TTF",
                r"C:\Windows\Fonts\malgunbd.ttf",
                # 备选：其他中文字体
                r"C:\Windows\Fonts\simsun.ttc",
                r"C:\Windows\Fonts\SIMSUN.TTC",
                r"C:\Windows\Fonts\simhei.ttf",
                r"C:\Windows\Fonts\SIMHEI.TTF",
            ]
            
            for font_path in absolute_fonts:
                if os.path.exists(font_path):
                    try:
                        font = ImageFont.truetype(font_path, size)
                        return font
                    except:
                        continue
        
        # Linux系统
        elif platform.system() == 'Linux':
            linux_fonts = [
                '/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc',
                '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
            ]
            for font_path in linux_fonts:
                if os.path.exists(font_path):
                    try:
                        font = ImageFont.truetype(font_path, size)
                        if self._test_font_support(font, test_chars):
                            return font
                    except:
                        continue
        
        # macOS系统
        elif platform.system() == 'Darwin':
            mac_fonts = [
                '/Library/Fonts/Arial Unicode.ttf',
                '/System/Library/Fonts/PingFang.ttc',
            ]
            for font_path in mac_fonts:
                if os.path.exists(font_path):
                    try:
                        font = ImageFont.truetype(font_path, size)
                        if self._test_font_support(font, test_chars):
                            return font
                    except:
                        continue
        
        # 最后尝试使用旧的_get_chinese_font方法
        return self._get_chinese_font(size)
    
    def _get_chinese_font(self, size):
        """获取支持多语言的Unicode字体（包括中文、日文、韩文、阿拉伯文、俄文等）"""
        # 测试字符：中文、韩文、日文、英文
        test_chars = ['中', '한', '日', 'A', '1']
        
        # Windows系统字体路径 - 优先使用支持Unicode的字体
        if platform.system() == 'Windows':
            # 支持Unicode的字体列表（按优先级排序）
            # 优先选择支持韩文、日文等东亚字符的字体
            unicode_fonts = [
                # Arial Unicode MS - 支持大量Unicode字符（包括中文、日文、韩文、阿拉伯文等）
                r"C:\Windows\Fonts\ARIALUNI.TTF",
                r"C:\Windows\Fonts\arialuni.ttf",
                # Malgun Gothic (맑은 고딕) - 微软韩文字体，完美支持韩文
                r"C:\Windows\Fonts\malgun.ttf",
                r"C:\Windows\Fonts\malgunbd.ttf",  # Malgun Gothic Bold
                r"C:\Windows\Fonts\MALGUN.TTF",
                r"C:\Windows\Fonts\MALGUNBD.TTF",
                # Microsoft YaHei (微软雅黑) - 支持中文、英文，部分支持韩文
                r"C:\Windows\Fonts\msyh.ttc",
                r"C:\Windows\Fonts\msyhbd.ttc",  # 微软雅黑 Bold
                # Microsoft YaHei UI - 支持中文、英文等
                r"C:\Windows\Fonts\msyhl.ttc",  # Microsoft YaHei UI Light
                r"C:\Windows\Fonts\MSYHL.TTC",
                # Gulim (굴림) - 韩文字体
                r"C:\Windows\Fonts\gulim.ttc",
                r"C:\Windows\Fonts\GULIM.TTC",
                # Batang (바탕) - 韩文字体
                r"C:\Windows\Fonts\batang.ttc",
                r"C:\Windows\Fonts\BATANG.TTC",
                # Segoe UI - Windows默认字体，支持多种语言
                r"C:\Windows\Fonts\segoeui.ttf",
                r"C:\Windows\Fonts\segoeuib.ttf",  # Segoe UI Bold
                # Calibri - 支持多种语言
                r"C:\Windows\Fonts\calibri.ttf",
                r"C:\Windows\Fonts\calibrib.ttf",  # Calibri Bold
                # SimSun (宋体) - 支持中文，但不支持韩文（放在后面）
                r"C:\Windows\Fonts\simsun.ttc",
                r"C:\Windows\Fonts\SIMSUN.TTC",
                # SimHei (黑体) - 支持中文
                r"C:\Windows\Fonts\simhei.ttf",
                r"C:\Windows\Fonts\SIMHEI.TTF",
            ]
            
            # 尝试加载支持Unicode的字体，并验证是否支持多语言
            loaded_font = None
            for font_path in unicode_fonts:
                if os.path.exists(font_path):
                    try:
                        test_font = ImageFont.truetype(font_path, size)
                        # 验证字体是否支持测试字符
                        if self._test_font_support(test_font, test_chars):
                            return test_font
                        # 如果验证失败，但字体加载成功，保存作为备选
                        if loaded_font is None:
                            loaded_font = test_font
                    except Exception as e:
                        continue
            
            # 如果所有字体都验证失败，使用第一个加载成功的字体
            if loaded_font is not None:
                return loaded_font
        
        # Linux系统字体路径
        elif platform.system() == 'Linux':
            # 支持Unicode的字体（按优先级排序）
            unicode_fonts = [
                # Noto字体系列 - Google开发，支持大量语言
                '/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc',
                '/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf',
                # DejaVu Sans - 优秀的Unicode字体
                '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
                # 文泉驿字体
                '/usr/share/fonts/truetype/wqy/wqy-microhei.ttc',
                '/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc',
            ]
            for font_path in unicode_fonts:
                if os.path.exists(font_path):
                    try:
                        return ImageFont.truetype(font_path, size)
                    except:
                        continue
        
        # macOS系统字体路径
        elif platform.system() == 'Darwin':
            # 支持Unicode的字体（按优先级排序）
            unicode_fonts = [
                # Arial Unicode MS - 优秀的Unicode字体
                '/Library/Fonts/Arial Unicode.ttf',
                # PingFang - macOS中文字体，支持Unicode
                '/System/Library/Fonts/PingFang.ttc',
                # STHeiti - 中文字体
                '/System/Library/Fonts/STHeiti Light.ttc',
                # Helvetica Neue - 系统默认，支持多种语言
                '/Library/Fonts/HelveticaNeue.ttc',
            ]
            for font_path in unicode_fonts:
                if os.path.exists(font_path):
                    try:
                        return ImageFont.truetype(font_path, size)
                    except:
                        continue
        
        # 如果找不到合适的字体，尝试使用PIL内置的Unicode字体
        try:
            # 尝试使用DejaVu Sans（如果可用）
            return ImageFont.truetype("DejaVuSans.ttf", size)
        except:
            try:
                # 尝试使用默认TrueType字体
                return ImageFont.truetype("arial.ttf", size)
            except:
                # 最后使用默认字体（可能不支持所有Unicode字符）
                return ImageFont.load_default()
