#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Excel转换器
"""

import os
from pathlib import Path
import openpyxl
from openpyxl import Workbook, load_workbook
from conversion_core.core.base import Converter


class ExcelConverter(Converter):
    """Excel转换器"""
    
    def _extract_tables_improved(self, page):
        """改进的表格提取方法"""
        tables = []
        
        # 使用多种策略提取表格
        try:
            # 策略1: 使用pdfplumber的默认表格提取
            default_tables = page.extract_tables()
            if default_tables:
                tables.extend(default_tables)
            
            # 策略2: 调整表格设置，提高识别率
            table_settings = {
                "vertical_strategy": "lines_strict",
                "horizontal_strategy": "lines_strict",
                "min_words_vertical": 1,
                "min_words_horizontal": 1,
                "snap_tolerance": 3,
                "join_tolerance": 3,
                "edge_min_length": 3,
                "intersection_tolerance": 3
            }
            
            strict_tables = page.extract_tables(table_settings)
            if strict_tables:
                # 去重并合并
                for table in strict_tables:
                    if table not in tables:
                        tables.append(table)
            
            # 策略3: 宽松设置，捕获更多可能的表格
            loose_settings = {
                "vertical_strategy": "text",
                "horizontal_strategy": "text",
                "snap_tolerance": 5,
                "join_tolerance": 5
            }
            
            loose_tables = page.extract_tables(loose_settings)
            if loose_tables:
                for table in loose_tables:
                    if table not in tables and self._is_valid_table_data(table):
                        tables.append(table)
                        
        except Exception as e:
            print(f"[表格提取错误] {e}")
        
        # 清理和验证表格
        cleaned_tables = []
        for table in tables:
            cleaned_table = self._clean_table_data(table)
            if cleaned_table and len(cleaned_table) > 0:
                cleaned_tables.append(cleaned_table)
        
        return cleaned_tables
    
    def _is_valid_table_data(self, table):
        """验证表格数据是否有效"""
        if not table or len(table) < 2:
            return False
        
        # 检查是否有足够的非空单元格
        non_empty_cells = 0
        total_cells = 0
        
        for row in table:
            if row:
                for cell in row:
                    total_cells += 1
                    if cell and str(cell).strip():
                        non_empty_cells += 1
        
        # 至少50%的单元格有内容
        return total_cells > 0 and (non_empty_cells / total_cells) >= 0.3
    
    def _clean_table_data(self, table):
        """清理表格数据"""
        if not table:
            return None
        
        cleaned_table = []
        for row in table:
            if row:
                cleaned_row = []
                for cell in row:
                    if cell is not None:
                        # 清理单元格内容
                        cell_str = str(cell).strip()
                        # 移除多余的空白字符
                        cell_str = ' '.join(cell_str.split())
                        cleaned_row.append(cell_str)
                    else:
                        cleaned_row.append("")
                
                # 只保留非空行
                if any(cell for cell in cleaned_row):
                    cleaned_table.append(cleaned_row)
        
        return cleaned_table if cleaned_table else None
    
    def _extract_structured_text(self, page):
        """提取结构化文本"""
        try:
            text = page.extract_text()
            if not text:
                return []
            
            lines = text.split('\n')
            structured_lines = []
            
            for line in lines:
                line = line.strip()
                if line:
                    # 尝试检测是否为表格行（包含多个制表符或空格分隔的数据）
                    if '\t' in line or '  ' in line:
                        # 可能是表格数据，尝试分割
                        parts = []
                        if '\t' in line:
                            parts = [part.strip() for part in line.split('\t') if part.strip()]
                        else:
                            # 使用多个空格分割
                            parts = [part.strip() for part in line.split('  ') if part.strip()]
                        
                        if len(parts) > 1:
                            structured_lines.append(' | '.join(parts))
                        else:
                            structured_lines.append(line)
                    else:
                        structured_lines.append(line)
            
            return structured_lines
            
        except Exception as e:
            print(f"[文本提取错误] {e}")
            return []
    
    def _convert_pdf_to_excel_with_ocr(self, file_path: str, mode: str):
        """使用OCR将PDF转换为Excel，专注于表格提取"""
        try:
            from conversion_core.services.ocr_conversion_service import OCRConversionService
            from conversion_core.core.ocr_excel_generator import OCRExcelGenerator
            from conversion_core.services.api_key_manager import APIKeyManager
            
            # 获取API密钥
            api_key_manager = APIKeyManager()
            api_key = api_key_manager.get_api_key()
            
            if not api_key:
                return {
                    'success': False, 
                    'error': '未配置OCR API密钥，请在设置中配置百度千帆API密钥'
                }
            
            # 配置OCR参数，专注于表格识别
            ocr_config = {
                "model": "paddleocr-vl",  # 使用PaddleOCR-VL，对表格识别效果更好
                "has_formula": False,
                "has_seal": False,
                "has_chart": False,
                "target_format": "excel",
                "output_dir": self.output_path
            }
            
            # 创建OCR转换服务
            ocr_service = OCRConversionService(api_key, ocr_config)
            
            # 执行OCR转换
            def progress_wrapper(current_page, total_pages, message):
                # 将OCR进度映射到15%-85%
                progress = 15 + int((current_page / total_pages) * 70)
                self.update_progress(file_path, progress)
            
            ocr_result = ocr_service.convert(file_path, progress_wrapper)
            
            if not ocr_result.get("success"):
                return {
                    'success': False,
                    'error': f'OCR识别失败: {ocr_result.get("error", "未知错误")}'
                }
            
            # 使用OCR Excel生成器
            excel_config = {"mode": mode}
            generator = OCRExcelGenerator(self.output_path, excel_config)
            
            self.update_progress(file_path, 90)
            
            # 生成Excel文档
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            output_filename = f"{base_name}_OCR.xlsx"
            
            excel_result = generator.generate(ocr_result.get("pages", []), output_filename)
            
            self.update_progress(file_path, 100)
            
            if excel_result.get("success"):
                return {
                    'success': True,
                    'output_file': excel_result.get("output_file"),
                    'table_count': excel_result.get("table_count", 0),
                    'sheet_count': excel_result.get("sheet_count", 0)
                }
            else:
                return {
                    'success': False,
                    'error': f'Excel生成失败: {excel_result.get("error", "未知错误")}'
                }
                
        except ImportError as e:
            return {
                'success': False,
                'error': f'OCR功能不可用: {str(e)}'
            }
        except Exception as e:
            return {
                'success': False,
                'error': f'OCR转换失败: {str(e)}'
            }
    
    def convert(self, file_path):
        """转换文件为Excel"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.pdf':
                return self._convert_pdf_to_excel(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._convert_excel_to_excel(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._convert_ppt_to_excel(file_path)
            elif ext in ['.docx', '.doc']:
                return self._convert_word_to_excel(file_path)
            else:
                return {'success': False, 'error': f'不支持的文件格式: {ext}'}
        
        except Exception as e:
            return {'success': False, 'error': f'转换失败: {str(e)}'}
    
    def _convert_pdf_to_excel(self, file_path):
        """PDF转Excel"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        # 推送开始进度
        self.update_progress(file_path, 5)
        
        try:
            import pdfplumber
        except ImportError as e:
            return {'success': False, 'error': f'pdfplumber库导入失败: {str(e)}。请确保已正确安装pdfplumber库。'}
        
        mode = self.config.get('mode', 'per_page')
        use_ocr = self.config.get('ocr', False)
        
        self.update_progress(file_path, 15)
        
        # 如果启用OCR，使用OCR转换服务
        if use_ocr:
            return self._convert_pdf_to_excel_with_ocr(file_path, mode)
        
        try:
            with pdfplumber.open(file_path) as pdf:
                if mode == 'per_page':
                    # 每页转为1个工作表
                    wb = Workbook()
                    wb.remove(wb.active)  # 删除默认工作表
                    
                    total_pages = len(pdf.pages)
                    self.update_progress(file_path, 25)
                    
                    for i, page in enumerate(pdf.pages, 1):
                        ws = wb.create_sheet(title=f"第{i}页")
                        
                        # 改进的表格提取
                        tables = self._extract_tables_improved(page)
                        
                        if tables:
                            current_row = 1
                            for table_idx, table in enumerate(tables):
                                if table_idx > 0:
                                    current_row += 2  # 表格间空行
                                
                                for row_idx, row in enumerate(table):
                                    for col_idx, cell in enumerate(row):
                                        if cell is not None:
                                            # 清理单元格内容
                                            cell_value = str(cell).strip() if cell else ""
                                            ws.cell(row=current_row + row_idx, column=col_idx + 1, value=cell_value)
                                
                                current_row += len(table)
                        else:
                            # 如果没有表格，尝试智能文本提取
                            text_data = self._extract_structured_text(page)
                            for row_idx, line in enumerate(text_data, 1):
                                ws.cell(row=row_idx, column=1, value=line)
                        
                        # 更新进度 (25% ~ 85%)
                        progress = 25 + int((i / total_pages) * 60)
                        self.update_progress(file_path, progress)
                    
                    self.update_progress(file_path, 90)
                    output_file = self.get_output_filename(file_path, 'xlsx')
                    wb.save(output_file)
                    
                else:  # merge mode
                    # 多页合为1个工作表
                    wb = Workbook()
                    ws = wb.active
                    ws.title = "合并内容"
                    
                    total_pages = len(pdf.pages)
                    self.update_progress(file_path, 25)
                    
                    row_offset = 1
                    for page_idx, page in enumerate(pdf.pages, 1):
                        # 添加页面标识
                        if page_idx > 1:
                            ws.cell(row=row_offset, column=1, value=f"--- 第{page_idx}页 ---")
                            row_offset += 1
                        
                        # 使用改进的表格提取
                        tables = self._extract_tables_improved(page)
                        
                        if tables:
                            for table_idx, table in enumerate(tables):
                                if table_idx > 0:
                                    row_offset += 1  # 表格间空行
                                
                                for row in table:
                                    max_col = 1
                                    for col_idx, cell in enumerate(row):
                                        if cell is not None:
                                            cell_value = str(cell).strip() if cell else ""
                                            ws.cell(row=row_offset, column=col_idx + 1, value=cell_value)
                                            max_col = max(max_col, col_idx + 1)
                                    row_offset += 1
                        else:
                            # 使用结构化文本提取
                            text_data = self._extract_structured_text(page)
                            for line in text_data:
                                ws.cell(row=row_offset, column=1, value=line)
                                row_offset += 1
                        
                        row_offset += 1  # 页面之间空一行
                        
                        # 更新进度 (25% ~ 85%)
                        progress = 25 + int((page_idx / total_pages) * 60)
                        self.update_progress(file_path, progress)
                    
                    self.update_progress(file_path, 90)
                    output_file = self.get_output_filename(file_path, 'xlsx')
                    wb.save(output_file)
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_file}
        
        except FileNotFoundError:
            return {'success': False, 'error': 'PDF文件不存在或无法访问'}
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'PDF转Excel失败: {error_msg}'}
    
    def _convert_excel_to_excel(self, file_path):
        """Excel转Excel（主要是格式转换）"""
        try:
            # 读取原Excel文件
            source_wb = load_workbook(file_path, data_only=True)
            
            # 创建新工作簿
            wb = Workbook()
            wb.remove(wb.active)
            
            # 复制所有工作表
            for sheet_name in source_wb.sheetnames:
                source_ws = source_wb[sheet_name]
                target_ws = wb.create_sheet(title=sheet_name)
                
                # 复制所有单元格
                for row in source_ws.iter_rows(values_only=False):
                    for cell in row:
                        if cell.value is not None:
                            target_ws.cell(row=cell.row, column=cell.column, value=cell.value)
            
            output_file = self.get_output_filename(file_path, 'xlsx')
            wb.save(output_file)
            source_wb.close()
            
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            import traceback
            error_detail = traceback.format_exc()
            return {'success': False, 'error': f'Excel转换失败: {str(e)}\n详细信息: {error_detail[:200]}'}
    
    def _convert_ppt_to_excel(self, file_path):
        """PPT转Excel"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        input_ext = Path(file_path).suffix.lower()
        
        # python-pptx只支持.pptx格式，不支持.ppt格式（旧格式）
        if input_ext == '.ppt':
            return {'success': False, 'error': 'PPT格式（旧格式）不支持。python-pptx库仅支持PPTX格式（新格式）。请先将PPT文件转换为PPTX格式。'}
        
        try:
            from pptx import Presentation
            from pptx.exc import PackageNotFoundError
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}。请确保已正确安装python-pptx库。'}
        
        mode = self.config.get('mode', 'per_page')
        
        try:
            prs = Presentation(file_path)
            wb = Workbook()
            
            if mode == 'per_page':
                wb.remove(wb.active)
                
                for i, slide in enumerate(prs.slides, 1):
                    ws = wb.create_sheet(title=f"第{i}页")
                    
                    row = 1
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            ws.cell(row=row, column=1, value=shape.text)
                            row += 1
            else:  # merge mode
                ws = wb.active
                ws.title = "合并内容"
                
                row = 1
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            ws.cell(row=row, column=1, value=shape.text)
                            row += 1
                    row += 1  # 幻灯片之间空一行
            
            output_file = self.get_output_filename(file_path, 'xlsx')
            wb.save(output_file)
            
            return {'success': True, 'output_file': output_file}
        
        except PackageNotFoundError:
            return {'success': False, 'error': 'PPT文件格式错误或文件已损坏。请检查文件是否为有效的PPTX格式。'}
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'PPT转Excel失败: {error_msg}'}
    
    def _convert_word_to_excel(self, file_path):
        """Word转Excel"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.doc':
            return {'success': False, 'error': 'DOC格式（旧格式）暂不支持，请先转换为DOCX格式'}
        
        try:
            from docx import Document
            from docx.opc.exceptions import PackageNotFoundError
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}。请确保已正确安装python-docx库。'}
        
        mode = self.config.get('mode', 'per_page')
        
        try:
            doc = Document(file_path)
            wb = Workbook()
            
            if mode == 'per_page':
                wb.remove(wb.active)
                
                # 按段落分组，每个段落作为一个工作表（简化处理）
                # 或者将所有内容放到一个工作表
                ws = wb.create_sheet(title="Word内容")
                
                row = 1
                for para in doc.paragraphs:
                    if para.text.strip():
                        ws.cell(row=row, column=1, value=para.text)
                        row += 1
                
                # 处理表格
                for table_idx, table in enumerate(doc.tables, 1):
                    # 表格前空一行
                    row += 1
                    table_start_row = row
                    
                    for table_row_idx, row_obj in enumerate(table.rows):
                        for col_idx, cell in enumerate(row_obj.cells):
                            ws.cell(row=row, column=col_idx + 1, value=cell.text.strip())
                        row += 1
                    
                    # 表格后空一行
                    row += 1
            else:  # merge mode
                ws = wb.active
                ws.title = "合并内容"
                
                row = 1
                for para in doc.paragraphs:
                    if para.text.strip():
                        ws.cell(row=row, column=1, value=para.text)
                        row += 1
                
                # 处理表格
                for table in doc.tables:
                    row += 1  # 表格前空一行
                    for row_obj in table.rows:
                        for col_idx, cell in enumerate(row_obj.cells):
                            ws.cell(row=row, column=col_idx + 1, value=cell.text.strip())
                        row += 1
                    row += 1  # 表格后空一行
            
            output_file = self.get_output_filename(file_path, 'xlsx')
            wb.save(output_file)
            
            return {'success': True, 'output_file': output_file}
        
        except PackageNotFoundError:
            return {'success': False, 'error': 'Word文件格式错误或文件已损坏。请检查文件是否为有效的DOCX格式。'}
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'Word转Excel失败: {error_msg}'}
    
    def _convert_image_to_excel(self, file_path):
        """图片转Excel"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        try:
            from PIL import Image
        except ImportError as e:
            return {'success': False, 'error': f'Pillow库导入失败: {str(e)}'}
        
        try:
            # 打开图片获取信息
            img = Image.open(file_path)
            width, height = img.size
            
            # 创建Excel工作簿
            wb = Workbook()
            ws = wb.active
            ws.title = "图片信息"
            
            # 写入图片信息
            ws.cell(row=1, column=1, value="图片路径")
            ws.cell(row=1, column=2, value=file_path)
            
            ws.cell(row=2, column=1, value="宽度")
            ws.cell(row=2, column=2, value=width)
            
            ws.cell(row=3, column=1, value="高度")
            ws.cell(row=3, column=2, value=height)
            
            ws.cell(row=4, column=1, value="格式")
            ws.cell(row=4, column=2, value=img.format)
            
            # 注意：openpyxl不支持直接插入图片，只能保存图片信息
            # 如果需要插入图片，需要使用其他方法（如openpyxl-image-loader）
            
            output_file = self.get_output_filename(file_path, 'xlsx')
            wb.save(output_file)
            
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'图片转Excel失败: {error_msg}'}
