#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT转换器
"""

import os
from pathlib import Path
from conversion_core.core.base import Converter


class PPTConverter(Converter):
    """PPT转换器"""
    
    def convert(self, file_path):
        """转换文件为PPT"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.pdf':
                return self._convert_pdf_to_ppt(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._convert_excel_to_ppt(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._convert_ppt_to_ppt(file_path)
            elif ext in ['.docx', '.doc']:
                return self._convert_word_to_ppt(file_path)
            elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff', '.webp']:
                return self._convert_image_to_ppt(file_path)
            else:
                return {'success': False, 'error': f'不支持的文件格式: {ext}'}
        
        except Exception as e:
            return {'success': False, 'error': f'转换失败: {str(e)}'}
    
    def _convert_pdf_to_ppt(self, file_path):
        """PDF转PPT"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        # 推送开始进度
        self.update_progress(file_path, 5)
        
        try:
            import fitz  # PyMuPDF
            pymupdf_available = True
            print(f"[PDF转PPT] PyMuPDF可用，版本: {fitz.version}")
        except ImportError as e:
            pymupdf_available = False
            print(f"[PDF转PPT] PyMuPDF不可用: {e}，将使用pdfplumber（仅文本模式）")
        
        try:
            import pdfplumber
        except ImportError as e:
            if not pymupdf_available:
                return {'success': False, 'error': f'pdfplumber库导入失败: {str(e)}'}
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        try:
            from PIL import Image
        except ImportError as e:
            return {'success': False, 'error': f'Pillow库导入失败: {str(e)}'}
        
        import tempfile
        import io
        
        self.update_progress(file_path, 15)
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # 优先使用PyMuPDF（支持图片提取）
            if pymupdf_available:
                doc = fitz.open(file_path)
                total_pages = len(doc)
                
                if total_pages == 0:
                    doc.close()
                    return {'success': False, 'error': 'PDF文件为空，没有可转换的页面'}
                
                self.update_progress(file_path, 25)
                
                for page_idx in range(total_pages):
                    page = doc[page_idx]
                    
                    # 创建新幻灯片
                    slide_layout = prs.slide_layouts[6]  # 空白布局
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 提取页面为图片（高质量渲染）
                    # 这样可以保留PDF中的所有内容包括图片
                    mat = fitz.Matrix(2.0, 2.0)  # 2x缩放以提高质量
                    pix = page.get_pixmap(matrix=mat)
                    
                    # 保存为临时图片
                    img_data = pix.tobytes("png")
                    
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                        tmp_path = tmp_file.name
                        tmp_file.write(img_data)
                    
                    try:
                        # 计算图片在幻灯片中的位置和大小
                        # 保持宽高比，适应幻灯片
                        img = Image.open(tmp_path)
                        img_width, img_height = img.size
                        img.close()
                        
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height
                        
                        # 计算缩放比例
                        margin = Inches(0.3)
                        available_width = slide_width - 2 * margin
                        available_height = slide_height - 2 * margin
                        
                        width_ratio = available_width / Emu(img_width * 914400 / 96)  # 96 DPI
                        height_ratio = available_height / Emu(img_height * 914400 / 96)
                        scale = min(width_ratio, height_ratio, 1.0)  # 不放大，只缩小
                        
                        pic_width = Emu(img_width * 914400 / 96 * scale)
                        pic_height = Emu(img_height * 914400 / 96 * scale)
                        
                        # 居中放置
                        left = (slide_width - pic_width) / 2
                        top = (slide_height - pic_height) / 2
                        
                        slide.shapes.add_picture(tmp_path, int(left), int(top), int(pic_width), int(pic_height))
                    finally:
                        try:
                            os.unlink(tmp_path)
                        except:
                            pass
                    
                    # 更新进度 (25% ~ 85%)
                    progress = 25 + int(((page_idx + 1) / total_pages) * 60)
                    self.update_progress(file_path, progress)
                
                doc.close()
            else:
                # 回退到pdfplumber（仅文本）
                with pdfplumber.open(file_path) as pdf:
                    if len(pdf.pages) == 0:
                        return {'success': False, 'error': 'PDF文件为空，没有可转换的页面'}
                    
                    total_pages = len(pdf.pages)
                    self.update_progress(file_path, 25)
                    
                    for page_idx, page in enumerate(pdf.pages, 1):
                        # 创建新幻灯片
                        slide_layout = prs.slide_layouts[6]  # 空白布局
                        slide = prs.slides.add_slide(slide_layout)
                        
                        # 提取文本
                        text = page.extract_text()
                        
                        if text:
                            # 添加文本框 - 调整大小和字体以避免溢出
                            left = Inches(0.3)
                            top = Inches(0.3)
                            width = Inches(9.4)
                            height = Inches(6.9)
                            
                            text_box = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = text_box.text_frame
                            text_frame.word_wrap = True
                            
                            # 根据文本长度动态调整字体大小
                            text_length = len(text)
                            if text_length > 2000:
                                font_size = Pt(8)
                            elif text_length > 1000:
                                font_size = Pt(9)
                            elif text_length > 500:
                                font_size = Pt(10)
                            else:
                                font_size = Pt(11)
                            
                            # 添加文本段落
                            p = text_frame.paragraphs[0]
                            p.text = text
                            p.font.size = font_size
                        
                        # 更新进度 (25% ~ 85%)
                        progress = 25 + int((page_idx / total_pages) * 60)
                        self.update_progress(file_path, progress)
            
            # 确保至少有一张幻灯片
            if len(prs.slides) == 0:
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
            
            self.update_progress(file_path, 90)
            output_file = self.get_output_filename(file_path, 'pptx')
            prs.save(output_file)
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'PDF转PPT失败: {error_msg}'}
    
    def _convert_excel_to_ppt(self, file_path):
        """Excel转PPT - 改进版，使用表格而不是文本"""
        try:
            from openpyxl import load_workbook
        except ImportError as e:
            return {'success': False, 'error': f'openpyxl库导入失败: {str(e)}'}
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        try:
            wb = load_workbook(file_path, data_only=True)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # 获取实际使用的行列范围
                max_row = ws.max_row
                max_col = ws.max_column
                
                # 跳过空工作表
                if max_row == 0 or max_col == 0:
                    continue
                
                # 限制表格大小（避免太大放不下）
                max_display_rows = 15
                max_display_cols = 8
                
                actual_rows = min(max_row, max_display_rows)
                actual_cols = min(max_col, max_display_cols)
                
                # 创建新幻灯片（空白布局）
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 添加标题
                title_left = Inches(0.5)
                title_top = Inches(0.3)
                title_width = Inches(9)
                title_height = Inches(0.6)
                
                title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = sheet_name
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                
                # 添加表格
                table_left = Inches(0.5)
                table_top = Inches(1.2)
                table_width = Inches(9)
                table_height = Inches(5.5)
                
                # 创建表格
                table = slide.shapes.add_table(
                    actual_rows, 
                    actual_cols, 
                    table_left, 
                    table_top, 
                    table_width, 
                    table_height
                ).table
                
                # 填充表格数据
                for row_idx in range(actual_rows):
                    for col_idx in range(actual_cols):
                        cell_value = ws.cell(row=row_idx + 1, column=col_idx + 1).value
                        cell_text = str(cell_value) if cell_value is not None else ''
                        
                        # 限制单元格文本长度
                        if len(cell_text) > 50:
                            cell_text = cell_text[:47] + '...'
                        
                        table.cell(row_idx, col_idx).text = cell_text
                        
                        # 设置字体大小
                        cell = table.cell(row_idx, col_idx)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(10)
                        
                        # 第一行加粗（作为表头）
                        if row_idx == 0:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.bold = True
                
                # 如果数据被截断，添加提示
                if max_row > max_display_rows or max_col > max_display_cols:
                    note_top = Inches(7)
                    note_box = slide.shapes.add_textbox(Inches(0.5), note_top, Inches(9), Inches(0.3))
                    note_frame = note_box.text_frame
                    note_para = note_frame.paragraphs[0]
                    note_para.text = f"注: 仅显示前 {actual_rows} 行 × {actual_cols} 列数据"
                    note_para.font.size = Pt(10)
                    note_para.font.italic = True
            
            # 确保至少有一张幻灯片
            if len(prs.slides) == 0:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "Excel文件为空"
                title_para.font.size = Pt(24)
            
            output_file = self.get_output_filename(file_path, 'pptx')
            prs.save(output_file)
            
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            import traceback
            error_detail = traceback.format_exc()
            return {'success': False, 'error': f'Excel转PPT失败: {str(e)}\n详细信息: {error_detail[:300]}'}
    
    def _convert_ppt_to_ppt(self, file_path):
        """PPT转PPT（主要是格式转换）"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        input_ext = Path(file_path).suffix.lower()
        
        # python-pptx只支持.pptx格式，不支持.ppt格式（旧格式）
        if input_ext == '.ppt':
            return {'success': False, 'error': 'PPT格式（旧格式）不支持。python-pptx库仅支持PPTX格式（新格式）。请先将PPT文件转换为PPTX格式，或使用Microsoft PowerPoint打开并另存为PPTX格式。'}
        
        try:
            from pptx import Presentation
            from pptx.exc import PackageNotFoundError
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        try:
            # 直接复制文件（或可以进行格式转换）
            output_file = self.get_output_filename(file_path, 'pptx')
            
            # 读取并保存
            prs = Presentation(file_path)
            prs.save(output_file)
            
            return {'success': True, 'output_file': output_file}
        
        except PackageNotFoundError:
            return {'success': False, 'error': 'PPT文件格式错误或文件已损坏。请检查文件是否为有效的PPTX格式。'}
        except Exception as e:
            error_msg = str(e)
            # 截断过长的错误信息
            if len(error_msg) > 200:
                error_msg = error_msg[:200] + "..."
            return {'success': False, 'error': f'PPT转换失败: {error_msg}'}
    
    def _convert_word_to_ppt(self, file_path):
        """Word转PPT"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.doc':
            return {'success': False, 'error': 'DOC格式（旧格式）暂不支持，请先转换为DOCX格式'}
        
        try:
            from docx import Document
            from docx.opc.exceptions import PackageNotFoundError
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        use_ocr = self.config.get('ocr', False)
        
        try:
            doc = Document(file_path)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # 按段落创建幻灯片
            for para in doc.paragraphs:
                if para.text.strip():
                    # 使用标题和内容布局
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 设置标题（使用段落的前几个词）
                    title = slide.shapes.title
                    title_text = para.text.strip()[:50]  # 限制标题长度
                    title.text = title_text
                    
                    # 设置内容
                    content = slide.placeholders[1]
                    text_frame = content.text_frame
                    text_frame.word_wrap = True
                    text_frame.text = para.text.strip()
            
            # 如果没有任何段落，创建一张空幻灯片
            if len(prs.slides) == 0:
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
            
            output_file = self.get_output_filename(file_path, 'pptx')
            prs.save(output_file)
            
            return {'success': True, 'output_file': output_file}
        
        except PackageNotFoundError:
            return {'success': False, 'error': 'Word文件格式错误或文件已损坏。请检查文件是否为有效的DOCX格式。'}
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'Word转PPT失败: {error_msg}'}
    
    def _convert_image_to_ppt(self, file_path):
        """图片转PPT"""
        # 检查文件是否存在且不为空
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return {'success': False, 'error': '文件为空，无法转换'}
        
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        try:
            from PIL import Image
        except ImportError as e:
            return {'success': False, 'error': f'Pillow库导入失败: {str(e)}'}
        
        try:
            # 打开图片
            img = Image.open(file_path)
            
            # 创建PPT演示文稿
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # 创建空白幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 计算图片尺寸（保持宽高比，最大宽度9英寸）
            max_width = Inches(9)
            width, height = img.size
            aspect_ratio = height / width
            pic_width = max_width
            pic_height = Inches(aspect_ratio * 9)
            
            # 居中放置图片
            left = (prs.slide_width - pic_width) / 2
            top = (prs.slide_height - pic_height) / 2
            
            # 添加图片到幻灯片
            # 使用临时文件保存图片
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                img.save(tmp_path, 'PNG')
            
            try:
                slide.shapes.add_picture(tmp_path, left, top, pic_width, pic_height)
            finally:
                # 删除临时文件
                try:
                    os.unlink(tmp_path)
                except:
                    pass
            
            output_file = self.get_output_filename(file_path, 'pptx')
            prs.save(output_file)
            
            return {'success': True, 'output_file': output_file}
        
        except Exception as e:
            error_msg = str(e)
            if len(error_msg) > 150:
                error_msg = error_msg[:150] + "..."
            return {'success': False, 'error': f'图片转PPT失败: {error_msg}'}
